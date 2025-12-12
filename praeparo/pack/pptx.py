"""PPTX assembly helpers for pack outputs.

Slides without a PPTX template are deliberately skipped: PNG artefacts remain
available for downstream consumers, and PPTX composition stays an optional
layer on top of the pack PNG pipeline.
"""

# pyright: reportUnknownMemberType=false, reportUnknownVariableType=false, reportAttributeAccessIssue=false, reportGeneralTypeIssues=false

from __future__ import annotations

import copy
import io
import logging
from dataclasses import dataclass
from pathlib import Path
from typing import Iterable, Mapping, Sequence

from jinja2 import Environment

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE, PP_PLACEHOLDER
from pptx.shapes.group import GroupShape
import pptx.shapes.picture as pptx_picture
from pptx.text.text import _Run

from praeparo.models import PackConfig, PackSlide
from praeparo.pack.templating import create_pack_jinja_env
from praeparo.visuals.dax.planner_core import slugify


def _slug_for_slide(slide: PackSlide, index: int) -> str:
    if slide.id:
        return slugify(slide.id)
    if slide.title:
        return slugify(slide.title)
    return f"slide_{index}"

logger = logging.getLogger(__name__)


@dataclass(frozen=True)
class PlaceholderSize:
    width_px: int
    height_px: int


def _emu_to_px(value) -> int:
    """Convert English Metric Units to pixels assuming 96 DPI."""

    # 1 inch = 914400 EMUs, 96 DPI
    return int(round((int(value) / 914400) * 96))


def _clear_slides(presentation: Presentation) -> None:
    """Remove all slides from *presentation* preserving layouts and theme."""
    # python-pptx does not expose public removal; prefer starting from a fresh
    # presentation when possible. This helper is retained for future use but
    # intentionally left empty to avoid relying on private attributes.
    if presentation.slides:
        return


def _clone_slide(presentation: Presentation, template_slide) -> object:
    """Clone *template_slide* into *presentation* preserving shapes."""

    new_slide = presentation.slides.add_slide(template_slide.slide_layout)

    # Drop any shapes added by the layout so the clone matches the template.
    for shape in list(new_slide.shapes):
        new_slide.shapes._spTree.remove(shape._element)  # type: ignore[attr-defined]

    for shape in template_slide.shapes:
        if isinstance(shape, pptx_picture.Picture):
            img = io.BytesIO(shape.image.blob)
            pic = new_slide.shapes.add_picture(
                image_file=img,
                left=shape.left,
                top=shape.top,
                width=shape.width,
                height=shape.height,
            )
            pic.name = getattr(shape, "name", None) or pic.name
            continue

        newel = copy.deepcopy(shape._element)  # type: ignore[attr-defined]
        new_slide.shapes._spTree.insert_element_before(newel, "p:extLst")  # type: ignore[attr-defined]

    return new_slide


def _picture_shapes(slide) -> list:
    shapes = []
    for shape in slide.shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            shapes.append(shape)
            continue
        if getattr(shape, "is_placeholder", False):
            try:
                if shape.placeholder_format.type == PP_PLACEHOLDER.PICTURE:
                    shapes.append(shape)
            except ValueError:
                # Some placeholders (e.g. missing type) may raise; skip them.
                continue
    return shapes


def _picture_placeholder_shapes(slide) -> list:
    shapes = []
    for shape in slide.shapes:
        if getattr(shape, "is_placeholder", False):
            try:
                if shape.placeholder_format.type == PP_PLACEHOLDER.PICTURE:
                    shapes.append(shape)
            except ValueError:
                continue
    return shapes


def _replace_picture(shape, image_path: Path) -> None:
    slide = shape.part.slide
    box_left, box_top, box_width, box_height = shape.left, shape.top, shape.width, shape.height
    name = getattr(shape, "name", None)
    shape._element.getparent().remove(shape._element)

    pic = slide.shapes.add_picture(str(image_path), box_left, box_top)

    img_width, img_height = pic.width, pic.height
    if not img_width or not img_height or not box_width or not box_height:
        pic.left = box_left
        pic.top = box_top
        pic.width = box_width
        pic.height = box_height
    else:
        scale = min(box_width / img_width, box_height / img_height)
        new_width = int(img_width * scale)
        new_height = int(img_height * scale)

        pic.width = new_width
        pic.height = new_height
        pic.left = box_left + (box_width - new_width) // 2
        pic.top = box_top + (box_height - new_height) // 2

    if name:
        pic.name = name


def _notes_template_tags(presentation: Presentation) -> dict[str, object]:
    """Return a mapping of TEMPLATE_TAG id -> slide object."""

    tags: dict[str, object] = {}
    for slide in presentation.slides:
        try:
            text = "".join(slide.notes_slide.notes_text_frame.text.split())
        except Exception:
            continue
        if "TEMPLATE_TAG=" not in text:
            continue
        token = text.split("TEMPLATE_TAG=", 1)[1].split()[0]
        if token:
            tags[token] = slide
    return tags


def resolve_template_geometry(template_path: Path) -> tuple[dict[str, PlaceholderSize], dict[tuple[str, str], PlaceholderSize]]:
    """Inspect a pack PPTX template and extract placeholder geometry.

    Returns a pair of dictionaries:
    - slide_sizes: template_tag -> PlaceholderSize for slides with a single
      picture placeholder.
    - placeholder_sizes: (template_tag, placeholder_name) -> PlaceholderSize
      for named picture shapes/placeholders on multi-visual templates.
    """

    slide_sizes: dict[str, PlaceholderSize] = {}
    placeholder_sizes: dict[tuple[str, str], PlaceholderSize] = {}

    presentation = Presentation(template_path)
    tags = _notes_template_tags(presentation)

    for tag, slide in tags.items():
        try:
            placeholders = _picture_placeholder_shapes(slide)

            # Single-visual case: exactly one picture placeholder.
            if len(placeholders) == 1:
                shape = placeholders[0]
                slide_sizes[tag] = PlaceholderSize(width_px=_emu_to_px(shape.width), height_px=_emu_to_px(shape.height))

            # Named placeholders and picture shapes.
            for shape in slide.shapes:
                is_picture = shape.shape_type == MSO_SHAPE_TYPE.PICTURE
                is_placeholder = getattr(shape, "is_placeholder", False)
                has_name = bool(getattr(shape, "name", None))
                if (is_picture or is_placeholder) and has_name:
                    placeholder_sizes[(tag, shape.name)] = PlaceholderSize(
                        width_px=_emu_to_px(shape.width),
                        height_px=_emu_to_px(shape.height),
                    )
        except Exception:
            logger.exception(
                "Failed to resolve geometry for template slide; continuing",
                extra={"template_tag": tag, "template_path": str(template_path)},
            )
            continue

    return slide_sizes, placeholder_sizes


def _delete_template_slides(presentation: Presentation) -> None:
    """Remove slides tagged as templates using the private slide id list."""

    xml_slides = presentation.slides._sldIdLst  # type: ignore[attr-defined]
    slides = list(presentation.slides)

    for index, slide in reversed(list(enumerate(slides))):
        try:
            text = slide.notes_slide.notes_text_frame.text
        except Exception:
            continue
        if text and "TEMPLATE_TAG=" in text:
            xml_slides.remove(xml_slides[index])


def _iter_shapes(shapes) -> Iterable[object]:
    """Yield every shape in *shapes*, descending into nested group shapes."""

    for shape in shapes:
        yield shape
        if isinstance(shape, GroupShape):
            yield from _iter_shapes(shape.shapes)


def _render_jinja_in_slide(
    slide,
    env: Environment,
    context: Mapping[str, object],
) -> None:
    """Render any Jinja placeholders found in text runs across the slide.

    The rendering is deliberately conservative: only runs containing ``{{`` are
    touched, runs outside a template remain unchanged, and multi-run templates
    are recombined before rendering to preserve formatting for surrounding text.
    """

    for shape in _iter_shapes(slide.shapes):
        if not getattr(shape, "has_text_frame", False):
            continue

        text_frame = shape.text_frame
        if text_frame is None:
            continue

        for paragraph in text_frame.paragraphs:
            runs: list[_Run] = list(paragraph.runs)
            i = 0
            while i < len(runs):
                text = runs[i].text
                if "{{" not in text:
                    i += 1
                    continue

                fragment = text
                j = i
                while "}}" not in fragment and j + 1 < len(runs):
                    j += 1
                    fragment += runs[j].text

                if "}}" not in fragment:
                    i = j + 1
                    continue

                rendered = env.from_string(fragment).render(**context)
                runs[i].text = rendered

                for k in range(i + 1, j + 1):
                    runs[k].text = ""

                i = j + 1


def assemble_pack_pptx(
    *,
    pack: PackConfig,
    results: Sequence[object] | None,
    slide_pngs: Mapping[str, Path],
    placeholder_pngs: Mapping[str, Mapping[str, Path]] | None = None,
    result_path: Path,
    template_path: Path | None = None,
    allow_missing_pngs: bool = False,
) -> None:
    """Build a PPTX from pack results using optional slide templates.

    Slides without a template are skipped rather than failing the run.

    When allow_missing_pngs is True, missing slide or placeholder PNG artefacts
    are logged and the corresponding template placeholders are left unchanged so
    partial runs can still assemble a deck.
    """

    placeholder_pngs = placeholder_pngs or {}

    jinja_env = create_pack_jinja_env()
    context_payload: dict[str, object] = {}
    if pack.context:
        context_payload.update(dict(pack.context))

    if template_path:
        presentation = Presentation(template_path)
        template_lookup: Mapping[str, object] = _notes_template_tags(presentation)
    else:
        presentation = Presentation()
        template_lookup = {}

    def _lookup_placeholder(slide_slug: str, placeholder_id: str) -> Path | None:
        placeholders = placeholder_pngs.get(slide_slug, {})
        key_slug = f"{slide_slug}__{slugify(placeholder_id)}"
        return placeholders.get(key_slug)

    for index, slide in enumerate(pack.slides, start=1):
        slide_slug = _slug_for_slide(slide, index)

        if slide.template:
            template_slide = template_lookup.get(slide.template)
            if template_slide is None:
                raise ValueError(f"Template '{slide.template}' not found for slide '{slide_slug}'")
        else:
            logger.info(
                "Skipping PPTX assembly for slide without template",
                extra={"slide": slide_slug},
            )
            continue

        cloned = _clone_slide(presentation, template_slide)

        # Title placeholder, if present.
        for shape in cloned.shapes:
            if getattr(shape, "is_placeholder", False):
                try:
                    if shape.placeholder_format.type in (PP_PLACEHOLDER.TITLE, PP_PLACEHOLDER.CENTER_TITLE):
                        shape.text = slide.title
                        break
                except ValueError:
                    continue

        slide_context: dict[str, object] = dict(context_payload)
        slide_context.update(
            {
                "title": slide.title,
                "slide_index": index,
                "slide_id": slide.id,
                "slide_slug": slide_slug,
            }
        )
        _render_jinja_in_slide(cloned, jinja_env, slide_context)

        has_placeholders = bool(slide.placeholders)
        has_visual = slide.visual is not None
        has_image = bool(slide.image)

        if has_placeholders:
            placeholders = slide.placeholders or {}
            for placeholder_id in placeholders:
                picture_shapes = [shape for shape in cloned.shapes if getattr(shape, "name", None) == placeholder_id]
                if not picture_shapes:
                    picture_shapes = _picture_shapes(cloned)
                    picture_shapes = [shape for shape in picture_shapes if getattr(shape, "name", None) == placeholder_id]

                if not picture_shapes:
                    raise ValueError(f"Placeholder '{placeholder_id}' not found on template '{slide.template}'")

                image_path = _lookup_placeholder(slide_slug, placeholder_id)
                if image_path is None:
                    if allow_missing_pngs:
                        logger.warning(
                            "Missing PNG for placeholder; leaving template unchanged",
                            extra={"slide": slide_slug, "placeholder": placeholder_id, "template": slide.template},
                        )
                        continue
                    raise ValueError(f"Missing PNG for placeholder '{placeholder_id}' on slide '{slide_slug}'")

                _replace_picture(picture_shapes[0], image_path)
        elif has_visual or has_image:
            picture_shapes = _picture_placeholder_shapes(cloned)
            if len(picture_shapes) != 1:
                raise ValueError(
                    f"Slide '{slide_slug}' uses single-visual shorthand but template '{slide.template}' "
                    f"has {len(picture_shapes)} picture picture-placeholders."
                )

            image_path = slide_pngs.get(slide_slug)
            if image_path is None:
                if allow_missing_pngs:
                    logger.warning(
                        "Missing PNG for slide; leaving template unchanged",
                        extra={"slide": slide_slug, "template": slide.template},
                    )
                    continue
                raise ValueError(f"Missing PNG for slide '{slide_slug}'")

            _replace_picture(picture_shapes[0], image_path)
        else:
            logger.info(
                "Rendering template-only slide without visual or placeholders",
                extra={"slide": slide_slug, "template": slide.template},
            )

    result_path.parent.mkdir(parents=True, exist_ok=True)
    if template_path:
        _delete_template_slides(presentation)
    presentation.save(result_path)
    logger.info(
        "Wrote PPTX",
        extra={"result_file": str(result_path), "slide_count": len(presentation.slides)},
    )


__all__ = ["assemble_pack_pptx", "PlaceholderSize", "resolve_template_geometry"]
