"""Power BI connectivity helpers."""

from __future__ import annotations

import asyncio
import os
import time
from dataclasses import dataclass
from io import BytesIO
from pathlib import Path
from typing import Any, Literal, Mapping, cast

import httpx

from .env import ensure_env_loaded


PowerBIExportFormatName = Literal["png", "pptx", "pdf"]
PBI_DEFAULT_EXPORT_FORMAT_ENV_VAR = "PRAEPARO_PBI_DEFAULT_EXPORT_FORMAT"
PBI_DEFAULT_GROUP_ID_ENV_VAR = "PRAEPARO_PBI_WORKSPACE_ID"
PBI_DEFAULT_REPORT_ID_ENV_VAR = "PRAEPARO_PBI_DEFAULT_REPORT_ID"
PBI_DEFAULT_STITCH_SLIDES_ENV_VAR = "PRAEPARO_PBI_DEFAULT_STITCH_SLIDES"
PBI_EXPORT_POLL_INTERVAL_ENV_VAR = "PRAEPARO_PBI_EXPORT_POLL_INTERVAL"
PBI_EXPORT_TIMEOUT_ENV_VAR = "PRAEPARO_PBI_EXPORT_TIMEOUT"


class PowerBIConfigurationError(RuntimeError):
    """Raised when required Power BI configuration is missing."""


class PowerBIAuthenticationError(RuntimeError):
    """Raised when acquiring an access token fails."""


class PowerBIQueryError(RuntimeError):
    """Raised when a DAX query execution fails."""


class PowerBIExportError(RuntimeError):
    """Raised when Power BI ExportToFile fails."""


@dataclass(frozen=True)
class PowerBIExportDefaults:
    """Runtime defaults that shape how Praeparo materialises Power BI exports."""

    format: PowerBIExportFormatName = "png"
    stitch_slides: bool = True
    poll_interval: float = 2.0
    timeout: float = 300.0

    @classmethod
    def from_env(
        cls,
        env: Mapping[str, str] | None = None,
    ) -> "PowerBIExportDefaults":
        """Resolve optional export defaults from `.env` or the live environment."""

        if env is None:
            ensure_env_loaded()
        env = env or os.environ

        return cls(
            format=_parse_export_format_env(
                PBI_DEFAULT_EXPORT_FORMAT_ENV_VAR,
                env.get(PBI_DEFAULT_EXPORT_FORMAT_ENV_VAR),
                default="png",
            ),
            stitch_slides=_parse_bool_env(
                PBI_DEFAULT_STITCH_SLIDES_ENV_VAR,
                env.get(PBI_DEFAULT_STITCH_SLIDES_ENV_VAR),
                default=True,
            ),
            poll_interval=_parse_positive_float_env(
                PBI_EXPORT_POLL_INTERVAL_ENV_VAR,
                env.get(PBI_EXPORT_POLL_INTERVAL_ENV_VAR),
                default=2.0,
            ),
            timeout=_parse_positive_float_env(
                PBI_EXPORT_TIMEOUT_ENV_VAR,
                env.get(PBI_EXPORT_TIMEOUT_ENV_VAR),
                default=300.0,
            ),
        )


@dataclass
class PowerBISettings:
    """Configuration required to authenticate with Power BI."""

    tenant_id: str
    client_id: str
    client_secret: str
    refresh_token: str
    scope: str = "https://analysis.windows.net/powerbi/api/.default"

    @classmethod
    def from_env(cls, env: Mapping[str, str] | None = None) -> "PowerBISettings":
        if env is None:
            ensure_env_loaded()
        env = env or os.environ
        try:
            tenant_id = env["PRAEPARO_PBI_TENANT_ID"]
            client_id = env["PRAEPARO_PBI_CLIENT_ID"]
            client_secret = env["PRAEPARO_PBI_CLIENT_SECRET"]
            refresh_token = env["PRAEPARO_PBI_REFRESH_TOKEN"]
        except KeyError as exc:
            raise PowerBIConfigurationError(
                "Missing Power BI configuration environment variables."
            ) from exc

        scope = env.get("PRAEPARO_PBI_SCOPE", "https://analysis.windows.net/powerbi/api/.default")
        return cls(
            tenant_id=tenant_id,
            client_id=client_id,
            client_secret=client_secret,
            refresh_token=refresh_token,
            scope=scope,
        )


class PowerBIClient:
    """Client for executing DAX queries against Power BI datasets."""

    def __init__(self, settings: PowerBISettings, *, timeout: float = 30.0) -> None:
        self._settings = settings
        self._timeout = timeout
        self._client = httpx.AsyncClient(timeout=timeout)
        self._access_token: str | None = None
        self._expires_at: float | None = None
        self._lock = asyncio.Lock()

    async def __aenter__(self) -> "PowerBIClient":
        return self

    async def __aexit__(self, exc_type, exc, tb) -> None:
        await self.aclose()

    async def aclose(self) -> None:
        await self._client.aclose()

    async def get_access_token(self) -> str:
        async with self._lock:
            if self._access_token and self._expires_at and self._expires_at - 60 > time.time():
                return self._access_token

            token_url = (
                f"https://login.microsoftonline.com/{self._settings.tenant_id}/oauth2/v2.0/token"
            )
            data = {
                "client_id": self._settings.client_id,
                "client_secret": self._settings.client_secret,
                "grant_type": "refresh_token",
                "refresh_token": self._settings.refresh_token,
                "scope": self._settings.scope,
            }

            response = await self._client.post(token_url, data=data)
            if response.status_code != 200:
                raise PowerBIAuthenticationError(
                    f"Failed to acquire Power BI access token: {response.status_code} {response.text}"
                )

            payload = response.json()
            access_token = payload.get("access_token")
            if not access_token:
                raise PowerBIAuthenticationError("Access token missing in authentication response.")

            expires_in = payload.get("expires_in")
            self._access_token = access_token
            if isinstance(expires_in, (int, float)):
                self._expires_at = time.time() + float(expires_in)
            else:
                self._expires_at = None
            return access_token

    async def execute_dax(
        self,
        dataset_id: str,
        query: str,
        *,
        group_id: str | None = None,
    ) -> list[dict[str, Any]]:
        token = await self.get_access_token()

        base_url = "https://api.powerbi.com/v1.0/myorg"
        if group_id:
            url = f"{base_url}/groups/{group_id}/datasets/{dataset_id}/executeQueries"
        else:
            url = f"{base_url}/datasets/{dataset_id}/executeQueries"

        response = await self._client.post(
            url,
            headers={"Authorization": f"Bearer {token}"},
            json={"queries": [{"query": query}]},
            timeout=600
        )

        if response.status_code != 200:
            raise PowerBIQueryError(
                f"Power BI query execution failed: {response.status_code} {response.text}"
            )

        payload = response.json()
        try:
            tables = payload["results"][0]["tables"]
            rows = tables[0]["rows"]
        except (KeyError, IndexError, TypeError) as exc:
            raise PowerBIQueryError("Unexpected response shape from Power BI executeQueries.") from exc

        return [_normalise_row_keys(row) for row in rows]

    async def export_to_file(
        self,
        *,
        group_id: str,
        report_id: str,
        payload: Mapping[str, Any],
        dest_path: str | os.PathLike[str],
        mode: str = "report",
        poll_interval: float = 2.0,
        timeout: float = 300.0,
    ) -> str:
        """Call ExportToFile for reports or paginated reports and persist the result.

        The request kicks off an export job, polls until completion (or failure),
        then downloads the file to `dest_path`. Callers choose the export payload
        (format/pages/filters) so this helper stays transport-focused.
        """

        token = await self.get_access_token()
        base_url = _export_base_url(group_id, report_id, mode)
        headers = {"Authorization": f"Bearer {token}"}

        start = await self._client.post(f"{base_url}/ExportTo", headers=headers, json=payload)
        if start.status_code not in (200, 202):
            raise PowerBIExportError(
                f"ExportToFile failed ({start.status_code}): {start.text}"
            )

        try:
            export_id = start.json()["id"]
        except Exception as exc:
            raise PowerBIExportError("ExportToFile response missing export id.") from exc

        deadline = time.time() + timeout
        status = "Running"
        retry_after: float | None = None
        while status not in {"Succeeded", "Failed"}:
            # Pace polling using either Retry-After from the service or the configured interval.
            if time.time() > deadline:
                raise PowerBIExportError("ExportToFile polling timed out.")

            wait_for = retry_after if retry_after and retry_after > 0 else poll_interval
            await asyncio.sleep(wait_for)

            status_resp = await self._client.get(
                f"{base_url}/exports/{export_id}",
                headers=headers,
            )
            if status_resp.status_code not in (200, 202):
                raise PowerBIExportError(
                    f"Failed to poll export status ({status_resp.status_code}): {status_resp.text}"
                )
            payload = status_resp.json()
            status = payload.get("status", "Unknown")
            retry_after = _parse_retry_after(status_resp)

            if status == "Failed":
                message = payload.get("error", {}).get("message") if isinstance(payload, dict) else None
                raise PowerBIExportError(f"Export failed: {message or 'unknown error'}")

        file_resp = await self._client.get(
            f"{base_url}/exports/{export_id}/file",
            headers=headers,
        )
        if file_resp.status_code != 200:
            raise PowerBIExportError(
                f"Failed to download export file ({file_resp.status_code}): {file_resp.text}"
            )

        dest = os.fspath(dest_path)
        Path(dest).parent.mkdir(parents=True, exist_ok=True)
        Path(dest).write_bytes(file_resp.content)
        return dest


def extract_png_from_pptx_export(
    pptx_export_path: str | os.PathLike[str],
    *,
    dest_path: str | os.PathLike[str] | None = None,
    stitch_slides: bool = True,
) -> Path:
    """Extract a PNG sidecar from a PPTX export.

    Power BI frequently slices tall report pages across multiple PPTX slides.
    This helper mirrors the legacy Slick behaviour by selecting the dominant
    picture from each slide, applying crop metadata, and stitching the visible
    segments into one PNG.
    """

    source = Path(pptx_export_path)
    target = Path(dest_path) if dest_path is not None else source.with_suffix(".png")

    # Keep the extracted image beside the source export so callers can cache both.
    target.parent.mkdir(parents=True, exist_ok=True)
    target.write_bytes(_extract_png_blob_from_pptx_export(source, stitch_slides=stitch_slides))
    return target


def get_default_powerbi_report_id(env: Mapping[str, str] | None = None) -> str | None:
    """Return the optional default report id configured for Power BI visuals."""

    if env is None:
        ensure_env_loaded()
    env = env or os.environ

    raw = env.get(PBI_DEFAULT_REPORT_ID_ENV_VAR)
    if raw is None:
        return None

    value = raw.strip()
    return value or None


def get_default_powerbi_group_id(env: Mapping[str, str] | None = None) -> str | None:
    """Return the optional default workspace id configured for Power BI visuals."""

    if env is None:
        ensure_env_loaded()
    env = env or os.environ

    raw = env.get(PBI_DEFAULT_GROUP_ID_ENV_VAR)
    if raw is None:
        return None

    value = raw.strip()
    return value or None


__all__ = [
    "PowerBIExportDefaults",
    "PowerBIExportFormatName",
    "PowerBIClient",
    "PowerBISettings",
    "PowerBIAuthenticationError",
    "PowerBIConfigurationError",
    "PowerBIQueryError",
    "PowerBIExportError",
    "extract_png_from_pptx_export",
    "get_default_powerbi_group_id",
    "get_default_powerbi_report_id",
]


def _normalise_row_keys(row: dict[str, object]) -> dict[str, object]:
    normalised: dict[str, object] = {}
    for key, value in row.items():
        normalised[key] = value
        stripped = _strip_bracket_wrappers(key)
        if stripped and stripped not in normalised:
            normalised[stripped] = value
    return normalised


def _strip_bracket_wrappers(label: str) -> str | None:
    start = label.rfind("[")
    end = label.rfind("]")
    if start == -1 or end == -1 or end <= start + 1:
        return None
    candidate = label[start + 1 : end].strip()
    if not candidate or candidate == label:
        return None
    return candidate


def _export_base_url(group_id: str, report_id: str, mode: str) -> str:
    route = "rdlreports" if mode == "paginated" else "reports"
    return f"https://api.powerbi.com/v1.0/myorg/groups/{group_id}/{route}/{report_id}"


def _parse_retry_after(response) -> float | None:
    try:
        header = response.headers.get("Retry-After")
        if header is None:
            return None
        
        # Retry a bit faster
        return float(header) / 3
    except Exception:
        return None


def _parse_export_format_env(
    env_var: str,
    raw: str | None,
    *,
    default: PowerBIExportFormatName,
) -> PowerBIExportFormatName:
    if raw is None or not raw.strip():
        return default

    candidate = raw.strip().lower()
    if candidate in {"png", "pptx", "pdf"}:
        return cast(PowerBIExportFormatName, candidate)

    raise PowerBIConfigurationError(
        f"{env_var} must be one of: png, pptx, pdf."
    )


def _parse_bool_env(env_var: str, raw: str | None, *, default: bool) -> bool:
    if raw is None or not raw.strip():
        return default

    candidate = raw.strip().lower()
    if candidate in {"1", "true", "yes", "on"}:
        return True
    if candidate in {"0", "false", "no", "off"}:
        return False

    raise PowerBIConfigurationError(
        f"{env_var} must be a boolean value such as true/false or 1/0."
    )


def _parse_positive_float_env(env_var: str, raw: str | None, *, default: float) -> float:
    if raw is None or not raw.strip():
        return default

    try:
        value = float(raw.strip())
    except ValueError as exc:
        raise PowerBIConfigurationError(f"{env_var} must be a positive number.") from exc

    if value <= 0:
        raise PowerBIConfigurationError(f"{env_var} must be greater than zero.")
    return value


def _extract_png_blob_from_pptx_export(
    pptx_export_path: Path,
    *,
    stitch_slides: bool,
) -> bytes:
    """Read a PPTX export, then crop and optionally stitch its primary pictures."""

    from PIL import Image as PILImage
    from pptx import Presentation
    from pptx.shapes.picture import Picture as PptxPicture

    presentation = Presentation(pptx_export_path)

    # First collect the visible image segment from each slide so downstream logic
    # can either stitch them together or keep just the first one.
    segments: list[PILImage.Image] = []
    for index, slide in enumerate(presentation.slides):
        picture = _pick_main_picture(slide, picture_type=PptxPicture)
        if picture is None:
            continue

        try:
            segments.append(_load_segment_from_picture(picture, image_type=PILImage))
        except Exception as exc:
            raise PowerBIExportError(
                f"Failed to read picture data from PPTX slide {index}."
            ) from exc

    if not segments:
        raise PowerBIExportError("PPTX export did not contain any slide pictures to extract.")

    # When callers disable stitching, preserve the first visible segment exactly
    # as exported so they can inspect one slide without building a tall composite.
    if not stitch_slides or len(segments) == 1:
        return _encode_png_bytes(segments[0])

    stitched = _stitch_segments(segments, image_type=PILImage)
    return _encode_png_bytes(stitched)


def _pick_main_picture(slide, *, picture_type):
    pictures = [shape for shape in slide.shapes if isinstance(shape, picture_type)]
    if not pictures:
        return None
    return max(pictures, key=lambda picture: int(picture.width) * int(picture.height))


def _load_segment_from_picture(shape, *, image_type):
    """Load the picture blob and crop it to the portion PowerPoint actually shows."""

    base = image_type.open(BytesIO(shape.image.blob))
    if base.mode not in {"RGB", "RGBA"}:
        base = base.convert("RGBA")
    else:
        base = base.convert("RGBA")

    crop_left = float(getattr(shape, "crop_left", 0.0) or 0.0)
    crop_right = float(getattr(shape, "crop_right", 0.0) or 0.0)
    crop_top = float(getattr(shape, "crop_top", 0.0) or 0.0)
    crop_bottom = float(getattr(shape, "crop_bottom", 0.0) or 0.0)

    if crop_left == crop_right == crop_top == crop_bottom == 0.0:
        return base

    width, height = base.size
    left = max(0, min(width, int(round(width * crop_left))))
    right = max(0, min(width, int(round(width * (1.0 - crop_right)))))
    top = max(0, min(height, int(round(height * crop_top))))
    bottom = max(0, min(height, int(round(height * (1.0 - crop_bottom)))))

    if right <= left or bottom <= top:
        raise PowerBIExportError("PPTX crop metadata produced an invalid image rectangle.")

    return base.crop((left, top, right, bottom))


def _stitch_segments(segments, *, image_type):
    """Trim duplicate overlap between segments, then stack them vertically."""

    target_width = max(segment.width for segment in segments)
    normalised = [_normalise_segment_width(segment, width=target_width, image_type=image_type) for segment in segments]

    overlaps: list[int] = []
    for index in range(len(normalised) - 1):
        overlaps.append(_exact_overlap_rows(normalised[index], normalised[index + 1]))

    total_height = sum(segment.height for segment in normalised) - sum(overlaps)
    stitched = image_type.new("RGBA", (target_width, total_height))

    y_offset = 0
    for index, segment in enumerate(normalised):
        crop_top = overlaps[index - 1] if index > 0 else 0
        visible = segment.crop((0, crop_top, segment.width, segment.height)) if crop_top else segment
        stitched.paste(visible, (0, y_offset))
        y_offset += visible.height

    return stitched


def _normalise_segment_width(segment, *, width: int, image_type):
    if segment.width == width:
        return segment

    resized_height = int(round(segment.height * (width / segment.width)))
    return segment.resize((width, resized_height), image_type.Resampling.LANCZOS)


def _exact_overlap_rows(bottom_segment, top_segment, *, max_search: int = 150) -> int:
    """Return how many top rows from *top_segment* duplicate *bottom_segment*."""

    if bottom_segment.width != top_segment.width:
        return 0

    row_bytes = bottom_segment.width * 4
    max_overlap = min(max_search, bottom_segment.height, top_segment.height)
    bottom_bytes = bottom_segment.tobytes()
    top_bytes = top_segment.tobytes()

    for overlap in range(max_overlap, 0, -1):
        bottom_slice = bottom_bytes[
            (bottom_segment.height - overlap) * row_bytes : bottom_segment.height * row_bytes
        ]
        top_slice = top_bytes[: overlap * row_bytes]
        if bottom_slice == top_slice:
            return overlap
    return 0


def _encode_png_bytes(image) -> bytes:
    buffer = BytesIO()
    image.save(buffer, format="PNG")
    return buffer.getvalue()
