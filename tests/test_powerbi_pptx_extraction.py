from __future__ import annotations

from pathlib import Path

from PIL import Image
from pptx import Presentation
from pptx.util import Inches

from praeparo.powerbi import extract_png_from_pptx_export


def _write_rows_png(path: Path, rows: list[tuple[int, int, int, int]], *, width: int = 4) -> None:
    image = Image.new("RGBA", (width, len(rows)))
    for y, colour in enumerate(rows):
        for x in range(width):
            image.putpixel((x, y), colour)
    image.save(path, format="PNG")


def _build_pptx_with_pictures(
    dest_path: Path,
    pictures: list[Path],
    *,
    crop_top: float = 0.0,
) -> None:
    presentation = Presentation()

    for picture_path in pictures:
        slide = presentation.slides.add_slide(presentation.slide_layouts[6])
        picture = slide.shapes.add_picture(
            str(picture_path),
            Inches(0.5),
            Inches(0.5),
            width=Inches(2),
            height=Inches(2),
        )
        if crop_top:
            picture.crop_top = crop_top

    presentation.save(dest_path)


def test_extract_png_from_pptx_export_stitches_slide_segments(tmp_path: Path) -> None:
    first = tmp_path / "first.png"
    second = tmp_path / "second.png"
    pptx_path = tmp_path / "export.pptx"

    red = (255, 0, 0, 255)
    green = (0, 255, 0, 255)
    blue = (0, 0, 255, 255)
    yellow = (255, 255, 0, 255)
    black = (0, 0, 0, 255)
    white = (255, 255, 255, 255)

    _write_rows_png(first, [red, green, blue, yellow])
    _write_rows_png(second, [yellow, black, white])
    _build_pptx_with_pictures(pptx_path, [first, second])

    output_path = extract_png_from_pptx_export(pptx_path, stitch_slides=True)
    extracted = Image.open(output_path).convert("RGBA")

    assert extracted.size == (4, 6)
    assert extracted.getpixel((0, 0)) == red
    assert extracted.getpixel((0, 1)) == green
    assert extracted.getpixel((0, 2)) == blue
    assert extracted.getpixel((0, 3)) == yellow
    assert extracted.getpixel((0, 4)) == black
    assert extracted.getpixel((0, 5)) == white


def test_extract_png_from_pptx_export_respects_crop_metadata(tmp_path: Path) -> None:
    source = tmp_path / "cropped.png"
    pptx_path = tmp_path / "cropped_export.pptx"

    red = (255, 0, 0, 255)
    green = (0, 255, 0, 255)
    blue = (0, 0, 255, 255)
    yellow = (255, 255, 0, 255)

    _write_rows_png(source, [red, green, blue, yellow])
    _build_pptx_with_pictures(pptx_path, [source], crop_top=0.25)

    output_path = extract_png_from_pptx_export(
        pptx_path,
        dest_path=tmp_path / "cropped.png",
        stitch_slides=False,
    )
    extracted = Image.open(output_path).convert("RGBA")

    assert extracted.size == (4, 3)
    assert extracted.getpixel((0, 0)) == green
    assert extracted.getpixel((0, 1)) == blue
    assert extracted.getpixel((0, 2)) == yellow
