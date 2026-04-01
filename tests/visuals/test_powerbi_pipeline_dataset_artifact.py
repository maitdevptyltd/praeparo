from __future__ import annotations

import json
from pathlib import Path
from typing import Mapping

from PIL import Image
from pptx import Presentation
from pptx.util import Inches

from praeparo.models import PowerBIRenderOptions, PowerBISource, PowerBIVisualConfig
from praeparo.pipeline import ExecutionContext, OutputTarget, PipelineOptions, VisualPipeline

# Import the visual module so its pipeline definition is registered.
import praeparo.visuals.powerbi as powerbi_visuals
from praeparo.powerbi import PowerBISettings


def _write_picture_pptx(dest_path: Path) -> None:
    source_png = dest_path.with_name("source.png")
    Image.new("RGBA", (8, 8), (255, 0, 0, 255)).save(source_png, format="PNG")

    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    slide.shapes.add_picture(str(source_png), Inches(0.5), Inches(0.5), width=Inches(2), height=Inches(2))
    presentation.save(dest_path)


def test_powerbi_pipeline_emits_stable_dataset_manifest(tmp_path: Path, monkeypatch) -> None:
    artefact_dir = tmp_path / "artefacts"
    exports_dir = tmp_path / "exports"

    monkeypatch.setenv("PRAEPARO_PBI_DEFAULT_EXPORT_FORMAT", "png")

    class FakePowerBIClient:
        def __init__(self, settings: PowerBISettings, *, timeout: float = 30.0) -> None:
            self._settings = settings
            self._timeout = timeout

        async def __aenter__(self) -> "FakePowerBIClient":
            return self

        async def __aexit__(self, exc_type, exc, tb) -> None:
            return None

        async def export_to_file(
            self,
            *,
            group_id: str,
            report_id: str,
            payload: Mapping[str, object],
            dest_path: str | Path,
            mode: str = "report",
            poll_interval: float = 2.0,
            timeout: float = 300.0,
        ) -> str:
            path = Path(dest_path)
            path.parent.mkdir(parents=True, exist_ok=True)
            path.write_bytes(b"stub")
            return str(path)

    monkeypatch.setattr(powerbi_visuals, "PowerBIClient", FakePowerBIClient)
    monkeypatch.setattr(
        powerbi_visuals.PowerBISettings,
        "from_env",
        classmethod(lambda cls, env=None: PowerBISettings(tenant_id="t", client_id="c", client_secret="s", refresh_token="r")),
    )

    config = PowerBIVisualConfig(
        title="Matters On Hold",
        mode="report",
        source=PowerBISource(group_id="group", report_id="report", page="Page 1"),
        filters=["MatterType eq 'Residential'"],
        render=PowerBIRenderOptions(format="png"),
    )
    pipeline = VisualPipeline()
    context = ExecutionContext(
        config_path=tmp_path / "visual.yaml",
        project_root=tmp_path,
        options=PipelineOptions(
            artefact_dir=artefact_dir,
            outputs=[],
            metadata={"build_artifacts_dir": exports_dir},
        ),
    )

    result = pipeline.execute(config, context)

    assert result.dataset_path is not None
    assert result.dataset_path.name == "data.json"
    assert result.dataset_path.exists()

    payload = json.loads(result.dataset_path.read_text(encoding="utf-8"))
    assert isinstance(payload, dict)

    assert payload.get("format") == "png"
    assert payload.get("pptx_path") is None
    assert payload.get("image_path", "").endswith(".png")
    export_payload = payload.get("export_payload")
    assert isinstance(export_payload, dict)
    assert export_payload.get("format") == "PNG"

    report_configuration = export_payload.get("powerBIReportConfiguration")
    assert isinstance(report_configuration, dict)
    assert report_configuration.get("pages") == [{"pageName": "Page 1"}]
    assert report_configuration.get("reportLevelFilters") == [{"filter": "MatterType eq 'Residential'"}]

    # Ensure we never persist Power BI credentials in dataset manifests.
    dataset_text = result.dataset_path.read_text(encoding="utf-8")
    assert "client_secret" not in dataset_text
    assert "refresh_token" not in dataset_text


def test_powerbi_pipeline_extracts_png_sidecar_from_pptx_and_honours_env_defaults(
    tmp_path: Path,
    monkeypatch,
) -> None:
    artefact_dir = tmp_path / "artefacts"
    exports_dir = tmp_path / "exports"
    requested_png = tmp_path / "requested.png"

    monkeypatch.setenv("PRAEPARO_PBI_DEFAULT_EXPORT_FORMAT", "pptx")
    monkeypatch.setenv("PRAEPARO_PBI_DEFAULT_STITCH_SLIDES", "false")
    monkeypatch.setenv("PRAEPARO_PBI_EXPORT_POLL_INTERVAL", "0.25")
    monkeypatch.setenv("PRAEPARO_PBI_EXPORT_TIMEOUT", "12.5")

    captured: dict[str, object] = {}

    class FakePowerBIClient:
        def __init__(self, settings: PowerBISettings, *, timeout: float = 30.0) -> None:
            self._settings = settings
            self._timeout = timeout

        async def __aenter__(self) -> "FakePowerBIClient":
            return self

        async def __aexit__(self, exc_type, exc, tb) -> None:
            return None

        async def export_to_file(
            self,
            *,
            group_id: str,
            report_id: str,
            payload: Mapping[str, object],
            dest_path: str | Path,
            mode: str = "report",
            poll_interval: float = 2.0,
            timeout: float = 300.0,
        ) -> str:
            captured["payload"] = dict(payload)
            captured["poll_interval"] = poll_interval
            captured["timeout"] = timeout

            path = Path(dest_path)
            path.parent.mkdir(parents=True, exist_ok=True)
            _write_picture_pptx(path)
            return str(path)

    monkeypatch.setattr(powerbi_visuals, "PowerBIClient", FakePowerBIClient)
    monkeypatch.setattr(
        powerbi_visuals.PowerBISettings,
        "from_env",
        classmethod(lambda cls, env=None: PowerBISettings(tenant_id="t", client_id="c", client_secret="s", refresh_token="r")),
    )

    config = PowerBIVisualConfig(
        title="Discharges Dashboard",
        mode="report",
        source=PowerBISource(group_id="group", report_id="report", page="Page 1"),
    )
    pipeline = VisualPipeline()
    context = ExecutionContext(
        config_path=tmp_path / "visual.yaml",
        project_root=tmp_path,
        options=PipelineOptions(
            artefact_dir=artefact_dir,
            outputs=[OutputTarget.png(requested_png)],
            metadata={"build_artifacts_dir": exports_dir},
        ),
    )

    result = pipeline.execute(config, context)

    assert requested_png.exists()
    assert result.dataset_path is not None

    payload = json.loads(result.dataset_path.read_text(encoding="utf-8"))
    assert payload.get("format") == "pptx"
    assert payload.get("export_path", "").endswith(".pptx")
    assert payload.get("pptx_path", "").endswith(".pptx")
    assert payload.get("image_path", "").endswith(".png")
    assert Path(payload["image_path"]).exists()

    export_payload = captured.get("payload")
    assert isinstance(export_payload, dict)
    assert export_payload.get("format") == "PPTX"
    assert captured.get("poll_interval") == 0.25
    assert captured.get("timeout") == 12.5
