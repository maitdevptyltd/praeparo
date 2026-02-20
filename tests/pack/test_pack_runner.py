from __future__ import annotations

import base64
from dataclasses import dataclass
from datetime import date
import threading
import time
from pathlib import Path
from typing import Any, Dict, Mapping, Tuple, cast

import plotly.graph_objects as go
import pytest
from pydantic import ConfigDict, Field, ValidationError, model_validator
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE, PP_PLACEHOLDER  # type: ignore[attr-defined]
from pptx.util import Inches
from PIL import Image

from praeparo.models import BaseVisualConfig, PackConfig, PackContext, PackPlaceholder, PackSlide, PackVisualRef
from praeparo.pack.filters import merge_odata_filters
from praeparo.pack.loader import load_pack_config
from praeparo.pack import PackExecutionError
from praeparo.pack.runner import PackPowerBIFailure, restitch_pack_pptx, run_pack
from praeparo.pack.metric_context import dump_context_payload
from praeparo.pack.templating import create_pack_jinja_env, render_value
from praeparo.pipeline import PipelineOptions, VisualExecutionResult, VisualPipeline, build_default_query_planner_provider
from praeparo.pipeline.outputs import OutputKind, PipelineOutputArtifact
from praeparo.powerbi import PowerBIQueryError
from praeparo.visuals.context import resolve_dax_context
from praeparo.visuals.dax.planner_core import slugify
from praeparo.visuals import VisualContextModel, register_visual_type


def test_pack_loader_and_templating(tmp_path: Path) -> None:
    pack_path = tmp_path / "pack.yaml"
    pack_path.write_text(
        """
schema: ing-pack-draft-1
context:
  lender_id: 201
  month: "2025-10-01"
filters:
  lender: "dim_lender/LenderId eq {{ lender_id }}"
  dates: "{{ odata_months_back_range('dim_calendar/month', month, 3) }}"
define: "DEFINE VAR Lender = {{ lender_id }}"
slides:
  - title: Example
    visual:
      ref: visuals/example.yaml
""",
        encoding="utf-8",
    )

    pack = load_pack_config(pack_path)
    env = create_pack_jinja_env()
    context_payload = dump_context_payload(pack.context)
    rendered_filters = render_value(pack.filters, env=env, context=context_payload)
    rendered_define = render_value(pack.define, env=env, context=context_payload)

    assert rendered_filters["lender"] == "dim_lender/LenderId eq 201"
    assert rendered_filters["dates"] == "dim_calendar/month ge 2025-08-01 and dim_calendar/month le 2025-10-01"
    assert rendered_define == "DEFINE VAR Lender = 201"


def test_pack_base_context_renders_registry_defines_against_pack_context(tmp_path: Path) -> None:
    (tmp_path / "registry" / "metrics").mkdir(parents=True)

    context_root = tmp_path / "registry" / "context"
    context_root.mkdir(parents=True)
    (context_root / "business_time.yaml").write_text(
        "\n".join(
            [
                "context:",
                "  business_time:",
                "    work_start: \"09:00\"",
                "    work_end: \"17:00\"",
                "define:",
                "  get_business_hours: |",
                "    FUNCTION GetCustomerBusinessHours =",
                "      (start_datetime : DateTime, end_datetime : DateTime) =>",
                "      RETURN",
                "        GetBusinessHours(",
                "          start_datetime,",
                "          end_datetime,",
                "          \"{{ business_time.work_start }}\",",
                "          \"{{ business_time.work_end }}\"",
                "        )",
            ]
        )
        + "\n",
        encoding="utf-8",
    )

    pack_path = tmp_path / "registry" / "customers" / "foo" / "pack.yaml"
    pack_path.parent.mkdir(parents=True, exist_ok=True)
    pack_path.write_text(
        "\n".join(
            [
                "schema: test-pack",
                "context:",
                "  business_time:",
                "    work_start: \"08:00\"",
                "    work_end: \"18:00\"",
            ]
        )
        + "\n",
        encoding="utf-8",
    )

    pack = load_pack_config(pack_path)
    env = create_pack_jinja_env()

    from praeparo.pack.runner import _resolve_pack_base_context_payload

    payload = _resolve_pack_base_context_payload(
        pack_path=pack_path,
        metadata=None,
        pack_context_layer=dump_context_payload(pack.context),
        env=env,
    )

    _, define_blocks = resolve_dax_context(base=payload, calculate=None, define=None)
    assert any("\"08:00\"" in block for block in define_blocks)
    assert any("\"18:00\"" in block for block in define_blocks)
    assert all("\"09:00\"" not in block for block in define_blocks)


def test_pack_base_context_prefers_metadata_context_over_pack_defaults(tmp_path: Path) -> None:
    (tmp_path / "registry" / "metrics").mkdir(parents=True)

    context_root = tmp_path / "registry" / "context"
    context_root.mkdir(parents=True)
    (context_root / "business_time.yaml").write_text(
        "\n".join(
            [
                "context:",
                "  business_time:",
                "    work_start: \"09:00\"",
                "    work_end: \"17:00\"",
                "define:",
                "  get_business_hours: |",
                "    FUNCTION GetCustomerBusinessHours =",
                "      () =>",
                "        GetBusinessHours(",
                "          BLANK(),",
                "          BLANK(),",
                "          \"{{ business_time.work_start }}\",",
                "          \"{{ business_time.work_end }}\"",
                "        )",
            ]
        )
        + "\n",
        encoding="utf-8",
    )

    pack_path = tmp_path / "registry" / "customers" / "foo" / "pack.yaml"
    pack_path.parent.mkdir(parents=True, exist_ok=True)
    pack_path.write_text(
        "\n".join(
            [
                "schema: test-pack",
                "context:",
                "  business_time:",
                "    work_start: \"08:00\"",
                "    work_end: \"18:00\"",
            ]
        )
        + "\n",
        encoding="utf-8",
    )

    pack = load_pack_config(pack_path)
    env = create_pack_jinja_env()

    from praeparo.pack.runner import _resolve_pack_base_context_payload

    payload = _resolve_pack_base_context_payload(
        pack_path=pack_path,
        metadata={
            "metrics_root": tmp_path / "registry" / "metrics",
            "context": {"business_time": {"work_start": "07:00", "work_end": "19:00"}},
        },
        pack_context_layer=dump_context_payload(pack.context),
        env=env,
    )

    _, define_blocks = resolve_dax_context(base=payload, calculate=None, define=None)
    assert any("\"07:00\"" in block for block in define_blocks)
    assert any("\"19:00\"" in block for block in define_blocks)
    assert all("\"08:00\"" not in block for block in define_blocks)
    assert all("\"18:00\"" not in block for block in define_blocks)


def test_merge_odata_filters_supports_dict_list_and_string() -> None:
    dict_merged = merge_odata_filters({"a": "one", "b": "two"}, {"b": "local", "c": "three"})
    assert dict_merged == {"a": "one", "b": "local", "c": "three"}

    list_merged = merge_odata_filters(["alpha"], ["beta", "gamma"])
    assert list_merged == ["alpha", "beta", "gamma"]

    string_merged = merge_odata_filters("first", "second")
    assert string_merged == ["first", "second"]

    inherit_global = merge_odata_filters(["base"], None)
    assert inherit_global == ["base"]


def test_run_pack_uses_project_root_for_discovery(monkeypatch, tmp_path: Path) -> None:
    monkeypatch.chdir(tmp_path)
    (tmp_path / "registry" / "metrics").mkdir(parents=True)

    config_dir = tmp_path / "configs"
    config_dir.mkdir()
    pack_path = config_dir / "pack.yaml"
    pack_path.write_text("{}", encoding="utf-8")
    visual_path = config_dir / "visual.yaml"
    visual_path.write_text("type: dummy_pack_noctx\n", encoding="utf-8")

    pack = PackConfig(
        schema="test-pack",
        slides=[
            PackSlide(title="Dummy Slide", visual=PackVisualRef(ref="visual.yaml")),
        ],
    )

    captured: Dict[str, object] = {}

    def dummy_loader(path: Path, payload=None, stack=()):
        return BaseVisualConfig(type="dummy_pack_noctx")

    from praeparo.pipeline.registry import (
        SchemaArtifact,
        DatasetArtifact,
        RenderOutcome,
        VisualPipelineDefinition,
        register_visual_pipeline,
    )

    def schema_builder(pipeline, config, context):
        assert context.dataset_context is not None
        captured["project_root"] = context.project_root
        captured["metrics_root"] = context.dataset_context.metrics_root
        return SchemaArtifact(value={})

    def dataset_builder(pipeline, config, schema, context):
        return DatasetArtifact(value={}, filename="data.json")

    def renderer(pipeline, config, schema, dataset, context, outputs):
        return RenderOutcome(outputs=[])

    register_visual_type("dummy_pack_noctx", dummy_loader, overwrite=True, context_model=None)
    register_visual_pipeline(
        "dummy_pack_noctx",
        VisualPipelineDefinition(
            schema_builder=schema_builder,
            dataset_builder=dataset_builder,
            renderer=renderer,
        ),
        overwrite=True,
    )

    pipeline = VisualPipeline(planner_provider=build_default_query_planner_provider())
    results = run_pack(
        pack_path,
        pack,
        project_root=tmp_path,
        output_root=tmp_path / "artefacts",
        base_options=PipelineOptions(),
        pipeline=pipeline,
        env=create_pack_jinja_env(),
    )

    assert results
    assert captured["project_root"] == tmp_path.resolve()
    assert captured["metrics_root"] == (tmp_path / "registry" / "metrics").resolve()


def test_run_pack_resolves_inherited_visual_refs_relative_to_declaring_pack(tmp_path: Path) -> None:
    (tmp_path / "registry" / "metrics").mkdir(parents=True)

    base_dir = tmp_path / "registry" / "packs" / "base"
    child_dir = tmp_path / "registry" / "packs" / "configs"
    base_dir.mkdir(parents=True, exist_ok=True)
    child_dir.mkdir(parents=True, exist_ok=True)

    base_path = base_dir / "base.yaml"
    base_path.write_text(
        """
schema: base-pack
slides:
  - id: overview
    title: Overview
    visual:
      ref: ./visuals/example.yaml
""",
        encoding="utf-8",
    )

    child_path = child_dir / "child.yaml"
    child_path.write_text(
        """
schema: child-pack
extends: ../base/base.yaml
""",
        encoding="utf-8",
    )

    pack = load_pack_config(child_path)

    captured: dict[str, Path] = {}

    def visual_loader(path: Path, payload=None, stack=()):
        captured["visual_path"] = path
        return BaseVisualConfig(type="dummy_inherited_ref")

    from praeparo.pipeline.registry import (
        SchemaArtifact,
        DatasetArtifact,
        RenderOutcome,
        VisualPipelineDefinition,
        register_visual_pipeline,
    )

    register_visual_pipeline(
        "dummy_inherited_ref",
        VisualPipelineDefinition(
            schema_builder=lambda pipeline, config, context: SchemaArtifact(value={}),
            dataset_builder=lambda pipeline, config, schema, context: DatasetArtifact(value={}, filename="data.json"),
            renderer=lambda pipeline, config, schema, dataset, context, outputs: RenderOutcome(outputs=[]),
        ),
        overwrite=True,
    )

    run_pack(
        child_path,
        pack,
        project_root=tmp_path,
        output_root=tmp_path / "out",
        base_options=PipelineOptions(),
        visual_loader=visual_loader,
    )

    assert captured["visual_path"] == (base_dir / "visuals" / "example.yaml").resolve()


def test_run_pack_resolves_inherited_placeholder_images_relative_to_declaring_pack(tmp_path: Path) -> None:
    (tmp_path / "registry" / "metrics").mkdir(parents=True)

    base_dir = tmp_path / "registry" / "packs" / "base"
    child_dir = tmp_path / "registry" / "packs" / "configs"
    base_assets = base_dir / "assets"
    base_assets.mkdir(parents=True, exist_ok=True)
    child_dir.mkdir(parents=True, exist_ok=True)

    logo_path = base_assets / "logo.png"
    _write_coloured_png(logo_path, colour=(0, 255, 0, 255))

    base_path = base_dir / "base.yaml"
    base_path.write_text(
        """
schema: base-pack
slides:
  - id: branded
    title: Branded
    template: home
    placeholders:
      logo:
        image: ./assets/logo.png
""",
        encoding="utf-8",
    )

    child_path = child_dir / "child.yaml"
    child_path.write_text(
        """
schema: child-pack
extends: ../base/base.yaml
""",
        encoding="utf-8",
    )

    pack = load_pack_config(child_path)
    run_pack(
        child_path,
        pack,
        project_root=tmp_path,
        output_root=tmp_path / "out",
        base_options=PipelineOptions(),
    )

def test_run_pack_resolves_registry_anchored_visual_refs(tmp_path: Path) -> None:
    (tmp_path / "registry" / "metrics").mkdir(parents=True)

    pack_path = tmp_path / "registry" / "customers" / "foo" / "pack.yaml"
    pack_path.parent.mkdir(parents=True, exist_ok=True)
    pack_path.write_text("{}", encoding="utf-8")

    visual_path = tmp_path / "registry" / "visuals" / "powerbi" / "pbi.yaml"
    visual_path.parent.mkdir(parents=True, exist_ok=True)
    visual_path.write_text("type: matrix\n", encoding="utf-8")

    pack = PackConfig(
        schema="test-pack",
        slides=[
            PackSlide(title="Anchored ref", visual=PackVisualRef(ref='@/visuals/powerbi/pbi.yaml')),
        ],
    )

    captured: Dict[str, Path] = {}

    def stub_visual_loader(path: Path) -> BaseVisualConfig:
        captured["visual_path"] = path
        return BaseVisualConfig(type="matrix")

    pipeline = _StubPipeline()
    run_pack(
        pack_path,
        pack,
        project_root=tmp_path,
        output_root=tmp_path / "artefacts",
        base_options=PipelineOptions(),
        visual_loader=stub_visual_loader,
        pipeline=cast(Any, pipeline),
        env=create_pack_jinja_env(),
    )

    assert captured["visual_path"] == visual_path.resolve()


def test_run_pack_rejects_registry_anchored_visual_refs_that_escape(tmp_path: Path) -> None:
    (tmp_path / "registry" / "metrics").mkdir(parents=True)

    pack_path = tmp_path / "registry" / "customers" / "foo" / "pack.yaml"
    pack_path.parent.mkdir(parents=True, exist_ok=True)
    pack_path.write_text("{}", encoding="utf-8")

    pack = PackConfig(
        schema="test-pack",
        slides=[
            PackSlide(title="Escaping ref", visual=PackVisualRef(ref='@/../x.yaml')),
        ],
    )

    pipeline = _StubPipeline()
    with pytest.raises(PackExecutionError) as excinfo:
        run_pack(
            pack_path,
            pack,
            project_root=tmp_path,
            output_root=tmp_path / "artefacts",
            base_options=PipelineOptions(),
            visual_loader=lambda _: BaseVisualConfig(type="matrix"),
            pipeline=cast(Any, pipeline),
            env=create_pack_jinja_env(),
        )
    exc = excinfo.value
    assert isinstance(exc.__cause__, ValueError)
    assert "@/../x.yaml" in str(exc.__cause__)


def test_run_pack_applies_inline_overrides_to_referenced_slide_visual(tmp_path: Path) -> None:
    (tmp_path / "registry" / "metrics").mkdir(parents=True)

    pack_path = tmp_path / "pack.yaml"
    pack_path.write_text("{}", encoding="utf-8")

    referenced_visual = tmp_path / "one.yaml"
    referenced_visual.write_text("type: matrix\n", encoding="utf-8")

    pack = PackConfig(
        schema="test-pack",
        slides=[
            PackSlide(
                id="slide-1",
                title="Inline override",
                visual=PackVisualRef.model_validate({"ref": "one.yaml", "title": "Override"}),
            ),
        ],
    )

    pipeline = _StubPipeline()

    def stub_visual_loader(path: Path) -> BaseVisualConfig:
        assert path == referenced_visual.resolve()
        return BaseVisualConfig(type="matrix", title="Base")

    run_pack(
        pack_path,
        pack,
        project_root=tmp_path,
        output_root=tmp_path / "artefacts",
        base_options=PipelineOptions(),
        visual_loader=stub_visual_loader,
        pipeline=cast(Any, pipeline),
        env=create_pack_jinja_env(),
    )

    assert pipeline.calls
    assert pipeline.calls[0][0].title == "Override"


def test_run_pack_applies_inline_overrides_to_referenced_placeholder_visual(tmp_path: Path) -> None:
    (tmp_path / "registry" / "metrics").mkdir(parents=True)

    pack_path = tmp_path / "pack.yaml"
    pack_path.write_text("{}", encoding="utf-8")

    referenced_visual = tmp_path / "one.yaml"
    referenced_visual.write_text("type: matrix\n", encoding="utf-8")

    pack = PackConfig(
        schema="test-pack",
        slides=[
            PackSlide(
                id="slide-1",
                title="Inline override placeholder",
                placeholders={
                    "chart": PackPlaceholder(
                        visual=PackVisualRef.model_validate({"ref": "one.yaml", "title": "Override"})
                    ),
                },
            ),
        ],
    )

    pipeline = _StubPipeline()

    def stub_visual_loader(path: Path) -> BaseVisualConfig:
        assert path == referenced_visual.resolve()
        return BaseVisualConfig(type="matrix", title="Base")

    run_pack(
        pack_path,
        pack,
        project_root=tmp_path,
        output_root=tmp_path / "artefacts",
        base_options=PipelineOptions(),
        visual_loader=stub_visual_loader,
        pipeline=cast(Any, pipeline),
        env=create_pack_jinja_env(),
    )

    assert pipeline.calls
    assert pipeline.calls[0][0].title == "Override"


def test_run_pack_rejects_unknown_inline_overrides_for_referenced_visual_with_context(tmp_path: Path) -> None:
    (tmp_path / "registry" / "metrics").mkdir(parents=True)

    pack_path = tmp_path / "pack.yaml"
    pack_path.write_text("{}", encoding="utf-8")

    referenced_visual = tmp_path / "one.yaml"
    referenced_visual.write_text("type: matrix\n", encoding="utf-8")

    pack = PackConfig(
        schema="test-pack",
        slides=[
            PackSlide(
                id="slide-1",
                title="Inline override failure",
                visual=PackVisualRef.model_validate({"ref": "one.yaml", "not_a_field": 1}),
            ),
        ],
    )

    pipeline = _StubPipeline()

    def stub_visual_loader(path: Path) -> BaseVisualConfig:
        assert path == referenced_visual.resolve()
        return BaseVisualConfig(type="matrix", title="Base")

    with pytest.raises(PackExecutionError) as excinfo:
        run_pack(
            pack_path,
            pack,
            project_root=tmp_path,
            output_root=tmp_path / "artefacts",
            base_options=PipelineOptions(),
            visual_loader=stub_visual_loader,
            pipeline=cast(Any, pipeline),
            env=create_pack_jinja_env(),
        )

    message = str(excinfo.value)
    assert str(pack_path) in message
    assert "id=slide-1" in message
    assert "one.yaml" in message
    assert "not_a_field" in message


def test_run_pack_preserves_excluded_discriminator_when_applying_overrides(tmp_path: Path) -> None:
    """Regression: override re-validation must not drop an excluded `type` discriminator."""

    class _ExcludedTypeVisualConfig(BaseVisualConfig):
        model_config = ConfigDict(extra="forbid", populate_by_name=True)
        type: str | None = Field(default=None, exclude=True)
        title: str | None = None

    (tmp_path / "registry" / "metrics").mkdir(parents=True)

    pack_path = tmp_path / "pack.yaml"
    pack_path.write_text("{}", encoding="utf-8")

    referenced_visual = tmp_path / "one.yaml"
    referenced_visual.write_text("type: matrix\n", encoding="utf-8")

    pack = PackConfig(
        schema="test-pack",
        slides=[
            PackSlide(
                id="slide-1",
                title="Excluded type override",
                visual=PackVisualRef.model_validate({"ref": "one.yaml", "title": "Override"}),
            ),
        ],
    )

    pipeline = _StubPipeline()

    def stub_visual_loader(path: Path) -> BaseVisualConfig:
        assert path == referenced_visual.resolve()
        return _ExcludedTypeVisualConfig(type="matrix", title="Base")

    run_pack(
        pack_path,
        pack,
        project_root=tmp_path,
        output_root=tmp_path / "artefacts",
        base_options=PipelineOptions(),
        visual_loader=stub_visual_loader,
        pipeline=cast(Any, pipeline),
        env=create_pack_jinja_env(),
    )

    assert pipeline.calls
    executed = pipeline.calls[0][0]
    assert executed.type == "matrix"
    assert executed.title == "Override"


def test_run_pack_applies_series_add_to_referenced_visual(tmp_path: Path) -> None:
    class _SeriesVisualConfig(BaseVisualConfig):
        model_config = ConfigDict(extra="forbid", populate_by_name=True)
        type: str | None = Field(default="series_visual")
        title: str | None = None
        series: list[dict[str, object]] = Field(default_factory=list)

    (tmp_path / "registry" / "metrics").mkdir(parents=True)

    pack_path = tmp_path / "pack.yaml"
    pack_path.write_text("{}", encoding="utf-8")

    referenced_visual = tmp_path / "series.yaml"
    referenced_visual.write_text("type: matrix\n", encoding="utf-8")

    pack = PackConfig(
        schema="test-pack",
        slides=[
            PackSlide(
                id="slide-1",
                title="Series add",
                visual=PackVisualRef.model_validate(
                    {
                        "ref": "series.yaml",
                        "series_add": [
                            {
                                "id": "other_lender",
                                "label": "Other Lender",
                            }
                        ],
                    }
                ),
            ),
        ],
    )

    pipeline = _StubPipeline()

    def stub_visual_loader(path: Path) -> BaseVisualConfig:
        assert path == referenced_visual.resolve()
        return _SeriesVisualConfig(
            type="series_visual",
            title="Base",
            series=[{"id": "customer_lender", "label": "ORDE"}],
        )

    run_pack(
        pack_path,
        pack,
        project_root=tmp_path,
        output_root=tmp_path / "artefacts",
        base_options=PipelineOptions(),
        visual_loader=stub_visual_loader,
        pipeline=cast(Any, pipeline),
        env=create_pack_jinja_env(),
    )

    assert pipeline.calls
    executed = cast(Any, pipeline.calls[0][0])
    assert [entry["id"] for entry in executed.series] == ["customer_lender", "other_lender"]
    assert executed.series[1]["label"] == "Other Lender"


def test_run_pack_applies_series_remove_and_update_to_referenced_visual(tmp_path: Path) -> None:
    class _SeriesVisualConfig(BaseVisualConfig):
        model_config = ConfigDict(extra="forbid", populate_by_name=True)
        type: str | None = Field(default="series_visual")
        series: list[dict[str, object]] = Field(default_factory=list)

    (tmp_path / "registry" / "metrics").mkdir(parents=True)

    pack_path = tmp_path / "pack.yaml"
    pack_path.write_text("{}", encoding="utf-8")

    referenced_visual = tmp_path / "series.yaml"
    referenced_visual.write_text("type: matrix\n", encoding="utf-8")

    pack = PackConfig(
        schema="test-pack",
        slides=[
            PackSlide(
                id="slide-1",
                title="Series update/remove",
                visual=PackVisualRef.model_validate(
                    {
                        "ref": "series.yaml",
                        "series_remove": ["legacy_line"],
                        "series_update": [
                            {
                                "id": "customer_lender",
                                "patch": {
                                    "label": "ORDE",
                                    "style": {"color": "#5B9BD5"},
                                },
                            }
                        ],
                    }
                ),
            ),
        ],
    )

    pipeline = _StubPipeline()

    def stub_visual_loader(path: Path) -> BaseVisualConfig:
        assert path == referenced_visual.resolve()
        return _SeriesVisualConfig(
            type="series_visual",
            series=[
                {"id": "customer_lender", "label": "Customer"},
                {"id": "legacy_line", "label": "Legacy"},
            ],
        )

    run_pack(
        pack_path,
        pack,
        project_root=tmp_path,
        output_root=tmp_path / "artefacts",
        base_options=PipelineOptions(),
        visual_loader=stub_visual_loader,
        pipeline=cast(Any, pipeline),
        env=create_pack_jinja_env(),
    )

    assert pipeline.calls
    executed = cast(Any, pipeline.calls[0][0])
    assert [entry["id"] for entry in executed.series] == ["customer_lender"]
    assert executed.series[0]["label"] == "ORDE"
    assert executed.series[0]["style"]["color"] == "#5B9BD5"


def test_run_pack_rejects_unknown_series_update_id(tmp_path: Path) -> None:
    class _SeriesVisualConfig(BaseVisualConfig):
        model_config = ConfigDict(extra="forbid", populate_by_name=True)
        type: str | None = Field(default="series_visual")
        series: list[dict[str, object]] = Field(default_factory=list)

    (tmp_path / "registry" / "metrics").mkdir(parents=True)

    pack_path = tmp_path / "pack.yaml"
    pack_path.write_text("{}", encoding="utf-8")

    referenced_visual = tmp_path / "series.yaml"
    referenced_visual.write_text("type: matrix\n", encoding="utf-8")

    pack = PackConfig(
        schema="test-pack",
        slides=[
            PackSlide(
                id="slide-1",
                title="Series update invalid",
                visual=PackVisualRef.model_validate(
                    {
                        "ref": "series.yaml",
                        "series_update": [
                            {
                                "id": "missing_series",
                                "patch": {"label": "Updated"},
                            }
                        ],
                    }
                ),
            ),
        ],
    )

    pipeline = _StubPipeline()

    def stub_visual_loader(path: Path) -> BaseVisualConfig:
        assert path == referenced_visual.resolve()
        return _SeriesVisualConfig(
            type="series_visual",
            series=[{"id": "customer_lender", "label": "Customer"}],
        )

    with pytest.raises(PackExecutionError) as excinfo:
        run_pack(
            pack_path,
            pack,
            project_root=tmp_path,
            output_root=tmp_path / "artefacts",
            base_options=PipelineOptions(),
            visual_loader=stub_visual_loader,
            pipeline=cast(Any, pipeline),
            env=create_pack_jinja_env(),
        )

    message = str(excinfo.value)
    assert "series_update" in message
    cause = excinfo.value.__cause__
    assert cause is not None
    detailed = str(cause.__cause__ or cause)
    assert "missing_series" in detailed


class _StubPipeline:
    def __init__(self) -> None:
        self.calls: list[Tuple[BaseVisualConfig, PipelineOptions]] = []
        self._lock = threading.Lock()
        self.contexts: list = []

    def execute(self, visual: BaseVisualConfig, context) -> VisualExecutionResult:
        with self._lock:
            self.calls.append((visual, context.options))
            self.contexts.append(context)
        outputs = []
        for target in context.options.outputs:
            if visual.type == "frame":
                continue
            target.path.parent.mkdir(parents=True, exist_ok=True)
            target.path.write_text(visual.type, encoding="utf-8")
            outputs.append(PipelineOutputArtifact(kind=target.kind, path=target.path))
        return VisualExecutionResult(config=visual, outputs=outputs)


class _ConcurrentStubPipeline:
    def __init__(self, *, delay: float = 0.05, fail_case: str | None = None) -> None:
        self.calls: list[Tuple[BaseVisualConfig, PipelineOptions]] = []
        self.delay = delay
        self.fail_case = fail_case
        self.running_powerbi = 0
        self.max_running_powerbi = 0
        self._lock = threading.Lock()

    def execute(self, visual: BaseVisualConfig, context) -> VisualExecutionResult:
        is_powerbi = visual.type == "powerbi"
        if is_powerbi:
            with self._lock:
                self.running_powerbi += 1
                self.max_running_powerbi = max(self.max_running_powerbi, self.running_powerbi)
            try:
                if self.fail_case and context.case_key == self.fail_case:
                    raise RuntimeError("boom")
                if self.delay:
                    time.sleep(self.delay)
            finally:
                with self._lock:
                    self.running_powerbi -= 1

        path = context.options.outputs[0].path
        path.parent.mkdir(parents=True, exist_ok=True)
        path.write_text(visual.type, encoding="utf-8")

        outputs = [PipelineOutputArtifact(kind=OutputKind.PNG, path=path)]

        with self._lock:
            self.calls.append((visual, context.options))

        return VisualExecutionResult(config=visual, outputs=outputs)


def _write_png(path: Path) -> None:
    data = base64.b64decode(
        "iVBORw0KGgoAAAANSUhEUgAAAAEAAAABCAYAAAAfFcSJAAAADUlEQVR42mP8z8BQDwAE"
        "AAL/AXx5lXcAAAAASUVORK5CYII="
    )
    path.write_bytes(data)


def _write_coloured_png(path: Path, colour: tuple[int, int, int, int] = (255, 0, 0, 255)) -> None:
    Image.new("RGBA", (1, 1), colour).save(path, format="PNG")


def _picture_filenames_by_name(slide) -> dict[str, str]:
    names: dict[str, str] = {}
    for shape in slide.shapes:
        if shape.shape_type != MSO_SHAPE_TYPE.PICTURE:  # type: ignore[attr-defined]
            continue
        name = getattr(shape, "name", None)
        if not name:
            continue
        try:
            names[name] = Path(shape.image.filename).name
        except Exception:
            continue
    return names


def _picture_blobs_by_name(slide) -> dict[str, bytes]:
    blobs: dict[str, bytes] = {}
    for shape in slide.shapes:
        if shape.shape_type != MSO_SHAPE_TYPE.PICTURE:  # type: ignore[attr-defined]
            continue
        name = getattr(shape, "name", None)
        if not name:
            continue
        try:
            blobs[name] = bytes(shape.image.blob)
        except Exception:
            continue
    return blobs


class _PngPipeline(_StubPipeline):
    def execute(self, visual: BaseVisualConfig, context) -> VisualExecutionResult:
        path = context.options.outputs[0].path
        path.parent.mkdir(parents=True, exist_ok=True)
        _write_png(path)
        outputs = [PipelineOutputArtifact(kind=OutputKind.PNG, path=path)]
        with self._lock:
            self.calls.append((visual, context.options))
            self.contexts.append(context)
        return VisualExecutionResult(config=visual, outputs=outputs)


def _build_pack_template(path: Path) -> None:
    prs = Presentation()
    blank = prs.slide_layouts[6]
    picture_with_caption = prs.slide_layouts[8]

    tmp_img = path.parent / "tmp.png"
    _write_png(tmp_img)

    single = prs.slides.add_slide(picture_with_caption)
    image_ph = next(ph for ph in single.placeholders if ph.placeholder_format.type == PP_PLACEHOLDER.PICTURE)
    image_ph.name = "image"
    single.notes_slide.notes_text_frame.text = "TEMPLATE_TAG=single_image"

    two_up = prs.slides.add_slide(blank)
    left = two_up.shapes.add_picture(str(tmp_img), Inches(0.5), Inches(1), width=Inches(3), height=Inches(2.5))
    left.name = "left_chart"
    right = two_up.shapes.add_picture(str(tmp_img), Inches(4), Inches(1), width=Inches(3), height=Inches(2.5))
    right.name = "right_chart"
    two_up.notes_slide.notes_text_frame.text = "TEMPLATE_TAG=two_up"

    prs.save(path)
    tmp_img.unlink(missing_ok=True)


def _build_two_up_template_with_text(path: Path) -> None:
    prs = Presentation()
    blank = prs.slide_layouts[6]

    tmp_img = path.parent / "tmp.png"
    _write_png(tmp_img)

    two_up = prs.slides.add_slide(blank)
    left = two_up.shapes.add_picture(str(tmp_img), Inches(0.5), Inches(1), width=Inches(3), height=Inches(2.5))
    left.name = "left_chart"
    right = two_up.shapes.add_picture(str(tmp_img), Inches(4), Inches(1), width=Inches(3), height=Inches(2.5))
    right.name = "right_chart"

    tb = two_up.shapes.add_textbox(Inches(0.5), Inches(0.2), width=Inches(6.5), height=Inches(0.5))
    tb.name = "display_date_text"
    tb.text_frame.text = "seed"

    two_up.notes_slide.notes_text_frame.text = "TEMPLATE_TAG=two_up"

    prs.save(path)
    tmp_img.unlink(missing_ok=True)


def _build_title_template(path: Path, template_tag: str = "title_only") -> None:
    prs = Presentation()
    title_slide = prs.slides.add_slide(prs.slide_layouts[0])
    title_slide.shapes.title.text = "Template Title"
    title_slide.notes_slide.notes_text_frame.text = f"TEMPLATE_TAG={template_tag}"
    prs.save(path)


def test_run_pack_attaches_geometry_to_slide_metadata(tmp_path: Path) -> None:
    pack_path = tmp_path / "pack.yaml"
    pack_path.write_text("", encoding="utf-8")

    template_path = pack_path.parent / "pack_template.pptx"
    _build_pack_template(template_path)

    pack = PackConfig(
        schema="test-pack",
        slides=[
            PackSlide(title="Deck Slide", id="deck-slide", template="single_image", visual=PackVisualRef(ref="one.yaml")),
        ],
    )

    visuals: Dict[str, BaseVisualConfig] = {"one.yaml": BaseVisualConfig(type="matrix")}

    def _loader(path: Path, payload: Mapping[str, object] | None = None, stack: tuple[Path, ...] = ()) -> BaseVisualConfig:
        return visuals[path.name]

    pipeline = _PngPipeline()
    run_pack(
        pack_path,
        pack,
        project_root=pack_path.parent,
        output_root=tmp_path / "artefacts",
        base_options=PipelineOptions(),
        visual_loader=_loader,
        pipeline=cast(VisualPipeline[Any], pipeline),
        env=create_pack_jinja_env(),
    )

    assert pipeline.contexts, "Pipeline should capture at least one context"
    slide_meta = pipeline.contexts[0].options.metadata
    assert isinstance(slide_meta.get("width"), int) and slide_meta["width"] > 0
    assert isinstance(slide_meta.get("height"), int) and slide_meta["height"] > 0


def test_run_pack_attaches_placeholder_geometry(tmp_path: Path) -> None:
    pack_path = tmp_path / "pack.yaml"
    pack_path.write_text("", encoding="utf-8")

    template_path = pack_path.parent / "pack_template.pptx"
    prs = Presentation()
    blank = prs.slide_layouts[6]

    tmp_img = tmp_path / "tmp.png"
    _write_png(tmp_img)

    two_up = prs.slides.add_slide(blank)
    left = two_up.shapes.add_picture(str(tmp_img), Inches(0.5), Inches(1), width=Inches(3), height=Inches(2))
    left.name = "left_chart"
    right = two_up.shapes.add_picture(str(tmp_img), Inches(4), Inches(1), width=Inches(2), height=Inches(3))
    right.name = "right_chart"
    two_up.notes_slide.notes_text_frame.text = "TEMPLATE_TAG=two_up"
    prs.save(template_path)
    tmp_img.unlink(missing_ok=True)

    pack = PackConfig(
        schema="test-pack",
        slides=[
            PackSlide(
                title="Two Up",
                id="two-up",
                template="two_up",
                placeholders={
                    "left_chart": PackPlaceholder(visual=PackVisualRef(ref="left.yaml")),
                    "right_chart": PackPlaceholder(visual=PackVisualRef(ref="right.yaml")),
                },
            )
        ],
    )

    visuals: Dict[str, BaseVisualConfig] = {
        "left.yaml": BaseVisualConfig(type="matrix"),
        "right.yaml": BaseVisualConfig(type="matrix"),
    }

    def _loader(path: Path, payload: Mapping[str, object] | None = None, stack: tuple[Path, ...] = ()) -> BaseVisualConfig:
        return visuals[path.name]

    pipeline = _PngPipeline()
    run_pack(
        pack_path,
        pack,
        project_root=pack_path.parent,
        output_root=tmp_path / "artefacts",
        base_options=PipelineOptions(),
        visual_loader=_loader,
        pipeline=cast(VisualPipeline[Any], pipeline),
        env=create_pack_jinja_env(),
    )

    meta_by_case = {context.case_key: context.options.metadata for context in pipeline.contexts}
    left_meta = next((meta for key, meta in meta_by_case.items() if "left_chart" in key), None)
    right_meta = next((meta for key, meta in meta_by_case.items() if "right_chart" in key), None)
    assert left_meta and right_meta
    assert left_meta["width"] != right_meta["width"] or left_meta["height"] != right_meta["height"]
    assert left_meta["width"] > 0 and right_meta["width"] > 0
    assert left_meta["height"] > 0 and right_meta["height"] > 0


def test_run_pack_respects_cli_width_height_override(tmp_path: Path) -> None:
    pack_path = tmp_path / "pack.yaml"
    pack_path.write_text("", encoding="utf-8")

    template_path = pack_path.parent / "pack_template.pptx"
    _build_pack_template(template_path)

    pack = PackConfig(
        schema="test-pack",
        slides=[PackSlide(title="Deck Slide", id="deck-slide", template="single_image", visual=PackVisualRef(ref="one.yaml"))],
    )

    visuals: Dict[str, BaseVisualConfig] = {"one.yaml": BaseVisualConfig(type="matrix")}

    def _loader(path: Path, payload: Mapping[str, object] | None = None, stack: tuple[Path, ...] = ()) -> BaseVisualConfig:
        return visuals[path.name]

    pipeline = _PngPipeline()
    overrides = PipelineOptions(metadata={"width": 1234, "height": 567})
    run_pack(
        pack_path,
        pack,
        project_root=pack_path.parent,
        output_root=tmp_path / "artefacts",
        base_options=overrides,
        visual_loader=_loader,
        pipeline=cast(VisualPipeline[Any], pipeline),
        env=create_pack_jinja_env(),
    )

    assert pipeline.contexts
    slide_meta = pipeline.contexts[0].options.metadata
    assert slide_meta["width"] == 1234
    assert slide_meta["height"] == 567


def test_run_pack_without_template_skips_geometry(tmp_path: Path) -> None:
    pack_path = tmp_path / "nested" / "packs" / "pack.yaml"
    pack_path.parent.mkdir(parents=True, exist_ok=True)
    pack_path.write_text("", encoding="utf-8")

    pack = PackConfig(
        schema="test-pack",
        slides=[PackSlide(title="Deck Slide", id="deck-slide", template="single_image", visual=PackVisualRef(ref="one.yaml"))],
    )

    visuals: Dict[str, BaseVisualConfig] = {"one.yaml": BaseVisualConfig(type="matrix")}

    def _loader(path: Path, payload: Mapping[str, object] | None = None, stack: tuple[Path, ...] = ()) -> BaseVisualConfig:
        return visuals[path.name]

    pipeline = _PngPipeline()
    run_pack(
        pack_path,
        pack,
        project_root=pack_path.parent,
        output_root=tmp_path / "artefacts",
        base_options=PipelineOptions(),
        visual_loader=_loader,
        pipeline=cast(VisualPipeline[Any], pipeline),
        env=create_pack_jinja_env(),
    )

    assert pipeline.contexts
    slide_meta = pipeline.contexts[0].options.metadata
    assert "width" not in slide_meta
    assert "height" not in slide_meta


def test_run_pack_templates_slide_metadata(tmp_path: Path) -> None:
    pack_path = tmp_path / "pack.yaml"
    pack_path.write_text("", encoding="utf-8")

    pack = PackConfig(
        schema="test-pack",
        context=PackContext.model_validate({"customer": "AMP"}),
        slides=[
            PackSlide(
                title="{{customer}} Dashboard",
                notes="Welcome, {{customer}}",
                visual=PackVisualRef(ref="one.yaml"),
            )
        ],
    )

    visuals: Dict[str, BaseVisualConfig] = {"one.yaml": BaseVisualConfig(type="matrix")}

    def _loader(path: Path, payload: Mapping[str, object] | None = None, stack: tuple[Path, ...] = ()) -> BaseVisualConfig:
        return visuals[path.name]

    pipeline = _StubPipeline()
    results = run_pack(
        pack_path,
        pack,
        project_root=pack_path.parent,
        output_root=tmp_path / "artefacts",
        base_options=PipelineOptions(),
        visual_loader=_loader,
        pipeline=cast(VisualPipeline[Any], pipeline),
        env=create_pack_jinja_env(),
    )

    slide = pack.slides[0]
    assert slide.title == "AMP Dashboard"
    assert slide.notes == "Welcome, AMP"

    png_paths = {result.png_path for result in results if result.png_path}
    assert (tmp_path / "artefacts" / "[01]_amp_dashboard.png") in png_paths


def test_run_pack_templates_slide_title_for_pptx(tmp_path: Path) -> None:
    pack_path = tmp_path / "pack.yaml"
    pack_path.write_text("", encoding="utf-8")

    template_path = tmp_path / "title_template.pptx"
    _build_title_template(template_path, template_tag="title_only")
    result_path = tmp_path / "deck" / "rendered_title.pptx"

    pack = PackConfig(
        schema="test-pack",
        context=PackContext.model_validate({"customer": "AMP"}),
        slides=[
            PackSlide(
                title="{{customer}} Dashboard",
                template="title_only",
            )
        ],
    )

    pipeline = _StubPipeline()
    base_options = PipelineOptions(metadata={"result_file": result_path, "pptx_template": template_path})
    run_pack(
        pack_path,
        pack,
        project_root=pack_path.parent,
        output_root=tmp_path / "artefacts",
        base_options=base_options,
        pipeline=cast(VisualPipeline[Any], pipeline),
        env=create_pack_jinja_env(),
    )

    deck = Presentation(result_path)
    assert deck.slides
    assert deck.slides[0].shapes.title.text == "AMP Dashboard"


def test_run_pack_names_outputs_with_ordinal_prefix(tmp_path: Path) -> None:
    pack_path = tmp_path / "pack.yaml"
    pack_path.write_text("", encoding="utf-8")

    pack = PackConfig(
        schema="test-pack",
        slides=[
            PackSlide(title="Alpha", id="alpha", visual=PackVisualRef(ref="one.yaml")),
            PackSlide(title="Beta Slide", visual=PackVisualRef(ref="two.yaml")),
        ],
    )

    visuals: Dict[str, BaseVisualConfig] = {
        "one.yaml": BaseVisualConfig(type="matrix"),
        "two.yaml": BaseVisualConfig(type="matrix"),
    }

    def _loader(path: Path, payload: Mapping[str, object] | None = None, stack: tuple[Path, ...] = ()) -> BaseVisualConfig:
        return visuals[path.name]

    pipeline = _PngPipeline()
    results = run_pack(
        pack_path,
        pack,
        project_root=pack_path.parent,
        output_root=tmp_path / "artefacts",
        base_options=PipelineOptions(),
        visual_loader=_loader,
        pipeline=cast(VisualPipeline[Any], pipeline),
        env=create_pack_jinja_env(),
    )

    png_paths = {result.png_path for result in results if result.png_path}
    assert (tmp_path / "artefacts" / "[01]_alpha.png") in png_paths
    assert (tmp_path / "artefacts" / "[02]_beta_slide.png") in png_paths

    slide_dirs = {call[1].artefact_dir for call in pipeline.calls}
    assert (tmp_path / "artefacts" / "[01]_alpha") in slide_dirs
    assert (tmp_path / "artefacts" / "[02]_beta_slide") in slide_dirs


def test_run_pack_routes_visuals_and_emits_pngs(tmp_path: Path) -> None:
    pack_path = tmp_path / "pack.yaml"
    pack_path.write_text("", encoding="utf-8")

    pack = PackConfig(
        schema="test-pack",
        context=PackContext.model_validate({"lender_id": 7, "month": "2025-11-01"}),
        define="DEFINE VAR Lender = {{ lender_id }}",
        calculate={"lender": "'dim_lender'[LenderId] = {{ lender_id }}"},
        filters={"lender": "dim_lender/LenderId eq {{ lender_id }}"},
        slides=[
            PackSlide(
                title="PowerBI Slide",
                id="pbi-slide",
                visual=PackVisualRef(
                    ref="powerbi.yaml",
                    filters={"local": "fact/Status eq 'Active'"},
                    calculate=["'dim_channel'[Name] = \"Broker\""],
                ),
            ),
            PackSlide(
                title="Matrix Visual",
                calculate=["'dim_channel'[Name] = \"Slide\""],
                visual=PackVisualRef(ref="matrix.yaml", calculate=["'dim_channel'[Name] = \"Direct\""]),
            ),
            PackSlide(
                title="Skip Frame",
                visual=PackVisualRef(ref="frame.yaml"),
            ),
        ],
    )

    visuals: Dict[str, BaseVisualConfig] = {
        "powerbi.yaml": BaseVisualConfig(type="powerbi"),
        "matrix.yaml": BaseVisualConfig(type="matrix"),
        "frame.yaml": BaseVisualConfig(type="frame"),
    }

    def _loader(path: Path, payload: Mapping[str, object] | None = None, stack: tuple[Path, ...] = ()) -> BaseVisualConfig:
        return visuals[path.name]

    pipeline = _StubPipeline()
    results = run_pack(
        pack_path,
        pack,
        project_root=pack_path.parent,
        output_root=tmp_path / "artefacts",
        base_options=PipelineOptions(),
        visual_loader=_loader,
        pipeline=cast(VisualPipeline[Any], pipeline),
        env=create_pack_jinja_env(),
    )

    png_paths = {result.png_path for result in results if result.png_path}
    assert (tmp_path / "artefacts" / f"[01]_{slugify('pbi-slide')}.png") in png_paths
    assert (tmp_path / "artefacts" / "[02]_matrix_visual.png") in png_paths

    # Frame visual produces no PNG but should not error.
    assert any(result.result.outputs == [] for result in results if result.result.config.type == "frame")

    metadata_by_type: Dict[str, Dict[str, object]] = {call[0].type: cast(Dict[str, object], call[1].metadata) for call in pipeline.calls}
    powerbi_metadata = metadata_by_type["powerbi"]
    assert "powerbi_filters" in powerbi_metadata
    assert "dim_lender/LenderId eq 7" in str(powerbi_metadata["powerbi_filters"])

    matrix_metadata = metadata_by_type["matrix"]
    context_meta = matrix_metadata.get("context")
    assert isinstance(context_meta, dict)
    calculate_payload = context_meta.get("calculate")
    assert isinstance(calculate_payload, list)
    assert calculate_payload[0].startswith("'dim_lender'[LenderId] = 7")
    assert calculate_payload[1] == "'dim_channel'[Name] = \"Slide\""
    assert calculate_payload[2] == "'dim_channel'[Name] = \"Direct\""
    assert len(calculate_payload) == 3
    define_payload = context_meta.get("define")
    assert isinstance(define_payload, list)
    assert define_payload == ["DEFINE VAR Lender = 7"]


def test_run_pack_populates_typed_dax_context(tmp_path: Path) -> None:
    register_visual_type(
        "contextual",
        lambda path, payload=None, stack=(): BaseVisualConfig(type="contextual"),
        overwrite=True,
        context_model=VisualContextModel,
    )

    pack_path = tmp_path / "pack.yaml"
    pack_path.write_text("", encoding="utf-8")

    pack = PackConfig(
        schema="test-pack",
        calculate=["'dim_channel'[Name] = \"Base\""],
        define="DEFINE MEASURE Test[Value] = 1",
        slides=[
            PackSlide(
                title="Context Slide",
                calculate=["'dim_channel'[Name] = \"Slide\""],
                visual=PackVisualRef(ref="contextual.yaml", calculate=["'dim_channel'[Name] = \"Direct\""]),
            ),
        ],
    )

    visuals: Dict[str, BaseVisualConfig] = {"contextual.yaml": BaseVisualConfig(type="contextual")}

    def _loader(path: Path, payload: Mapping[str, object] | None = None, stack: tuple[Path, ...] = ()) -> BaseVisualConfig:
        return visuals[path.name]

    pipeline = _PngPipeline()
    results = run_pack(
        pack_path,
        pack,
        project_root=pack_path.parent,
        output_root=tmp_path / "artefacts",
        base_options=PipelineOptions(metadata={"metrics_root": "registry/metrics"}),
        visual_loader=_loader,
        pipeline=cast(VisualPipeline[Any], pipeline),
        env=create_pack_jinja_env(),
    )

    assert results
    assert pipeline.contexts
    visual_ctx = pipeline.contexts[0].visual_context
    assert visual_ctx is not None
    assert visual_ctx.dax.calculate == (
        "'dim_channel'[Name] = \"Base\"",
        "'dim_channel'[Name] = \"Slide\"",
        "'dim_channel'[Name] = \"Direct\"",
    )
    assert visual_ctx.dax.define == ("DEFINE MEASURE Test[Value] = 1",)


def test_run_pack_forwards_pack_context_into_visual_context(tmp_path: Path) -> None:
    class _TestContextModel(VisualContextModel):
        reference_date: date | None = None
        trailing_months: int = 3

        @model_validator(mode="before")
        @classmethod
        def _from_context(cls, values: Mapping[str, object]) -> Mapping[str, object]:
            data = dict(values)
            ctx = data.get("context") or {}
            if isinstance(ctx, Mapping):
                if "reference_date" not in data and "month" in ctx:
                    data["reference_date"] = ctx["month"]
                if "trailing_months" not in data and "trailing_months" in ctx:
                    data["trailing_months"] = ctx["trailing_months"]
            return data

    register_visual_type(
        "contextual_month",
        lambda path, payload=None, stack=(): BaseVisualConfig(type="contextual_month"),
        overwrite=True,
        context_model=_TestContextModel,
    )

    pack_path = tmp_path / "pack.yaml"
    pack_path.write_text("", encoding="utf-8")

    pack = PackConfig(
        schema="test-pack",
        context=PackContext.model_validate({"month": "2025-10-01", "trailing_months": 3}),
        slides=[
            PackSlide(
                title="Context Slide",
                visual=PackVisualRef(ref="contextual_month.yaml"),
            ),
        ],
    )

    visuals: Dict[str, BaseVisualConfig] = {"contextual_month.yaml": BaseVisualConfig(type="contextual_month")}

    def _loader(path: Path, payload: Mapping[str, object] | None = None, stack: tuple[Path, ...] = ()) -> BaseVisualConfig:
        return visuals[path.name]

    pipeline = _PngPipeline()
    results = run_pack(
        pack_path,
        pack,
        project_root=pack_path.parent,
        output_root=tmp_path / "artefacts",
        base_options=PipelineOptions(metadata={"metrics_root": "registry/metrics"}),
        visual_loader=_loader,
        pipeline=cast(VisualPipeline[Any], pipeline),
        env=create_pack_jinja_env(),
    )

    assert results
    assert pipeline.contexts
    visual_ctx = pipeline.contexts[0].visual_context
    assert isinstance(visual_ctx, _TestContextModel)
    assert visual_ctx.reference_date == date(2025, 10, 1)
    assert visual_ctx.trailing_months == 3


def test_run_pack_metadata_context_overrides_pack_context_in_visual_context(tmp_path: Path) -> None:
    class _TestContextModel(VisualContextModel):
        reference_date: date | None = None
        trailing_months: int = 3

        @model_validator(mode="before")
        @classmethod
        def _from_context(cls, values: Mapping[str, object]) -> Mapping[str, object]:
            data = dict(values)
            ctx = data.get("context") or {}
            if isinstance(ctx, Mapping):
                if "reference_date" not in data and "month" in ctx:
                    data["reference_date"] = ctx["month"]
                if "trailing_months" not in data and "trailing_months" in ctx:
                    data["trailing_months"] = ctx["trailing_months"]
            return data

    register_visual_type(
        "contextual_month_override",
        lambda path, payload=None, stack=(): BaseVisualConfig(type="contextual_month_override"),
        overwrite=True,
        context_model=_TestContextModel,
    )

    pack_path = tmp_path / "pack.yaml"
    pack_path.write_text("", encoding="utf-8")

    pack = PackConfig(
        schema="test-pack",
        context=PackContext.model_validate({"month": "2025-10-01", "trailing_months": 3}),
        slides=[
            PackSlide(
                title="Context Slide",
                visual=PackVisualRef(ref="contextual_month_override.yaml"),
            ),
        ],
    )

    visuals: Dict[str, BaseVisualConfig] = {
        "contextual_month_override.yaml": BaseVisualConfig(type="contextual_month_override")
    }

    def _loader(path: Path, payload: Mapping[str, object] | None = None, stack: tuple[Path, ...] = ()) -> BaseVisualConfig:
        return visuals[path.name]

    pipeline = _PngPipeline()
    results = run_pack(
        pack_path,
        pack,
        project_root=pack_path.parent,
        output_root=tmp_path / "artefacts",
        base_options=PipelineOptions(
            metadata={
                "metrics_root": "registry/metrics",
                "context": {"month": "2025-11-01", "trailing_months": 6},
            }
        ),
        visual_loader=_loader,
        pipeline=cast(VisualPipeline[Any], pipeline),
        env=create_pack_jinja_env(),
    )

    assert results
    assert pipeline.contexts
    visual_ctx = pipeline.contexts[0].visual_context
    assert isinstance(visual_ctx, _TestContextModel)
    assert visual_ctx.reference_date == date(2025, 11, 1)
    assert visual_ctx.trailing_months == 6


def test_run_pack_named_calculate_overrides_global(tmp_path: Path) -> None:
    register_visual_type(
        "contextual_override",
        lambda path, payload=None, stack=(): BaseVisualConfig(type="contextual_override"),
        overwrite=True,
        context_model=VisualContextModel,
    )

    pack_path = tmp_path / "pack.yaml"
    pack_path.write_text("", encoding="utf-8")

    pack = PackConfig(
        schema="test-pack",
        context=PackContext.model_validate({"lender_id": 7}),
        calculate={
            "lender": "'dim_lender'[LenderId] = {{ lender_id }}",
            "channel": "'dim_channel'[Name] = \"Base\"",
        },
        slides=[
            PackSlide(
                title="Context Slide",
                visual=PackVisualRef(
                    ref="contextual.yaml",
                    calculate={
                        "lender": "'dim_lender'[LenderId] = 9",
                        "region": "'dim_region'[Name] = \"NSW\"",
                    },
                ),
            ),
        ],
    )

    visuals: Dict[str, BaseVisualConfig] = {"contextual.yaml": BaseVisualConfig(type="contextual_override")}

    def _loader(path: Path, payload: Mapping[str, object] | None = None, stack: tuple[Path, ...] = ()) -> BaseVisualConfig:
        return visuals[path.name]

    pipeline = _StubPipeline()
    run_pack(
        pack_path,
        pack,
        project_root=pack_path.parent,
        output_root=tmp_path / "artefacts",
        base_options=PipelineOptions(),
        visual_loader=_loader,
        pipeline=cast(VisualPipeline[Any], pipeline),
        env=create_pack_jinja_env(),
    )

    assert pipeline.contexts
    visual_ctx = pipeline.contexts[0].visual_context
    assert visual_ctx is not None
    assert visual_ctx.dax.calculate == (
        "'dim_lender'[LenderId] = 9",
        "'dim_channel'[Name] = \"Base\"",
        "'dim_region'[Name] = \"NSW\"",
    )


def test_run_pack_honours_only_slides(tmp_path: Path) -> None:
    pack_path = tmp_path / "pack.yaml"
    pack_path.write_text("", encoding="utf-8")

    slides = [
        PackSlide(
            title="Keep Me",
            id="keep-1",
            visual=PackVisualRef(ref="one.yaml"),
        ),
        PackSlide(
            title="Also Keep",
            visual=PackVisualRef(ref="two.yaml"),
        ),
        PackSlide(
            title="Skip Me",
            id="skip-this",
            visual=PackVisualRef(ref="three.yaml"),
        ),
        PackSlide(
            title="Slug Target",
            id="Slug-ID",
            visual=PackVisualRef(ref="four.yaml"),
        ),
    ]

    pack = PackConfig(
        schema="test-pack",
        slides=slides,
    )

    visuals: Dict[str, BaseVisualConfig] = {
        "one.yaml": BaseVisualConfig(type="powerbi"),
        "two.yaml": BaseVisualConfig(type="matrix"),
        "three.yaml": BaseVisualConfig(type="powerbi"),
        "four.yaml": BaseVisualConfig(type="powerbi"),
    }

    def _loader(path: Path, payload: Mapping[str, object] | None = None, stack: tuple[Path, ...] = ()) -> BaseVisualConfig:
        return visuals[path.name]

    pipeline = _StubPipeline()
    results = run_pack(
        pack_path,
        pack,
        project_root=pack_path.parent,
        output_root=tmp_path / "artefacts",
        base_options=PipelineOptions(),
        visual_loader=_loader,
        pipeline=cast(VisualPipeline[Any], pipeline),
        env=create_pack_jinja_env(),
        only_slides=["keep-1", "Also Keep", slugify("Slug Target")],
    )

    executed_titles = {result.slide.title for result in results}
    assert executed_titles == {"Keep Me", "Also Keep", "Slug Target"}
    called_visuals = {call[0].type for call in pipeline.calls}
    assert called_visuals == {"powerbi", "matrix"}

    expected_pngs = {
        tmp_path / "artefacts" / "[01]_keep_1.png",
        tmp_path / "artefacts" / "[02]_also_keep.png",
        tmp_path / "artefacts" / "[04]_slug_id.png",
    }
    emitted_pngs = {result.png_path for result in results if result.png_path}
    assert expected_pngs == emitted_pngs


def test_run_pack_only_slides_allows_missing_pngs(tmp_path: Path) -> None:
    pack_path = tmp_path / "pack.yaml"
    pack_path.write_text("", encoding="utf-8")

    template_path = tmp_path / "pack_template.pptx"
    _build_pack_template(template_path)
    result_path = tmp_path / "deck" / "partial_pack.pptx"

    pack = PackConfig(
        schema="test-pack",
        slides=[
            PackSlide(
                title="Single Visual",
                id="single",
                template="single_image",
                visual=PackVisualRef(ref="one.yaml"),
            ),
            PackSlide(
                title="Placeholders",
                id="placeholders",
                template="two_up",
                placeholders={
                    "left_chart": PackPlaceholder(visual=PackVisualRef(ref="left.yaml")),
                    "right_chart": PackPlaceholder(visual=PackVisualRef(ref="right.yaml")),
                },
            ),
            PackSlide(
                title="Other Slide",
                id="other",
                template="single_image",
                visual=PackVisualRef(ref="two.yaml"),
            ),
        ],
    )

    visuals: Dict[str, BaseVisualConfig] = {
        "one.yaml": BaseVisualConfig(type="matrix"),
        "left.yaml": BaseVisualConfig(type="matrix"),
        "right.yaml": BaseVisualConfig(type="matrix"),
        "two.yaml": BaseVisualConfig(type="matrix"),
    }

    def _loader(path: Path, payload: Mapping[str, object] | None = None, stack: tuple[Path, ...] = ()) -> BaseVisualConfig:
        return visuals[path.name]

    pipeline = _PngPipeline()
    base_options = PipelineOptions(metadata={"result_file": result_path, "pptx_template": template_path})

    run_pack(
        pack_path,
        pack,
        project_root=pack_path.parent,
        output_root=tmp_path / "artefacts",
        base_options=base_options,
        visual_loader=_loader,
        pipeline=cast(VisualPipeline[Any], pipeline),
        env=create_pack_jinja_env(),
        only_slides=["single"],
    )

    assert len(pipeline.calls) == 1, "Only the targeted slide should execute"
    assert result_path.exists(), "PPTX should be assembled even with missing PNGs for skipped slides"
    deck = Presentation(result_path)
    assert len(deck.slides) == 3


def test_run_pack_queues_powerbi_and_respects_concurrency(tmp_path: Path) -> None:
    pack_path = tmp_path / "pack.yaml"
    pack_path.write_text("", encoding="utf-8")

    pack = PackConfig(
        schema="test-pack",
        slides=[
            PackSlide(title="PBI One", id="pbi-one", visual=PackVisualRef(ref="pbi1.yaml")),
            PackSlide(title="Matrix Slide", visual=PackVisualRef(ref="matrix.yaml")),
            PackSlide(title="PBI Two", id="pbi-two", visual=PackVisualRef(ref="pbi2.yaml")),
        ],
    )

    visuals: Dict[str, BaseVisualConfig] = {
        "pbi1.yaml": BaseVisualConfig(type="powerbi"),
        "pbi2.yaml": BaseVisualConfig(type="powerbi"),
        "matrix.yaml": BaseVisualConfig(type="matrix"),
    }

    def _loader(path: Path, payload: Mapping[str, object] | None = None, stack: tuple[Path, ...] = ()) -> BaseVisualConfig:
        return visuals[path.name]

    pipeline = _ConcurrentStubPipeline(delay=0.05)
    results = run_pack(
        pack_path,
        pack,
        project_root=pack_path.parent,
        output_root=tmp_path / "artefacts",
        base_options=PipelineOptions(),
        visual_loader=_loader,
        pipeline=cast(VisualPipeline[Any], pipeline),
        env=create_pack_jinja_env(),
        max_powerbi_concurrency=2,
    )

    png_paths = {result.png_path for result in results if result.png_path}
    expected_pngs = {
        tmp_path / "artefacts" / f"[01]_{slugify('pbi-one')}.png",
        tmp_path / "artefacts" / f"[03]_{slugify('pbi-two')}.png",
        tmp_path / "artefacts" / "[02]_matrix_slide.png",
    }
    assert expected_pngs == png_paths
    assert pipeline.max_running_powerbi <= 2
    assert len([item for item in results if item.result.config.type == "powerbi"]) == 2


def test_run_pack_raises_when_powerbi_job_fails(tmp_path: Path) -> None:
    pack_path = tmp_path / "pack.yaml"
    pack_path.write_text("", encoding="utf-8")

    failing_slug = slugify("fail-pbi")
    pack = PackConfig(
        schema="test-pack",
        slides=[
            PackSlide(title="OK Slide", id="ok-pbi", visual=PackVisualRef(ref="pbi1.yaml")),
            PackSlide(title="Fail Slide", id="fail-pbi", visual=PackVisualRef(ref="pbi2.yaml")),
        ],
    )

    visuals: Dict[str, BaseVisualConfig] = {
        "pbi1.yaml": BaseVisualConfig(type="powerbi"),
        "pbi2.yaml": BaseVisualConfig(type="powerbi"),
    }

    def _loader(path: Path, payload: Mapping[str, object] | None = None, stack: tuple[Path, ...] = ()) -> BaseVisualConfig:
        return visuals[path.name]

    pipeline = _ConcurrentStubPipeline(delay=0.0, fail_case=failing_slug)

    with pytest.raises(PackPowerBIFailure) as excinfo:
        run_pack(
            pack_path,
            pack,
            project_root=pack_path.parent,
            output_root=tmp_path / "artefacts",
            base_options=PipelineOptions(),
            visual_loader=_loader,
            pipeline=cast(VisualPipeline[Any], pipeline),
            env=create_pack_jinja_env(),
            max_powerbi_concurrency=2,
        )

    message = str(excinfo.value)
    assert "Power BI slide(s) failed" in message
    assert failing_slug in message
    assert "PackExecutionError" in message
    ok_png = tmp_path / "artefacts" / f"[01]_{slugify('ok-pbi')}.png"
    assert ok_png.exists()


def test_run_pack_wraps_slide_failures_with_pack_execution_error(tmp_path: Path) -> None:
    (tmp_path / "registry" / "metrics").mkdir(parents=True)

    pack_path = tmp_path / "pack.yaml"
    pack_path.write_text("", encoding="utf-8")

    pack = PackConfig(
        schema="test-pack",
        slides=[
            PackSlide(title="Broken Slide", id="broken", visual=PackVisualRef(ref="visuals/broken.yaml")),
        ],
    )

    def _loader(path: Path, payload: Mapping[str, object] | None = None, stack: tuple[Path, ...] = ()) -> BaseVisualConfig:
        return BaseVisualConfig(type="matrix")

    class _FailingPipeline:
        def execute(self, visual: BaseVisualConfig, context):  # noqa: ANN001
            raise PowerBIQueryError("400 Bad Request: DAX execution failed")

    pipeline = _FailingPipeline()

    with pytest.raises(PackExecutionError) as excinfo:
        run_pack(
            pack_path,
            pack,
            project_root=pack_path.parent,
            output_root=tmp_path / "artefacts",
            base_options=PipelineOptions(),
            visual_loader=_loader,
            pipeline=cast(Any, pipeline),
            env=create_pack_jinja_env(),
        )

    exc = excinfo.value
    message = str(exc)
    assert str(pack_path) in message
    assert "phase=visual_execute" in message
    assert "Broken Slide" in message
    assert "visual_ref=visuals/broken.yaml" in message
    assert isinstance(exc.__cause__, PowerBIQueryError)


def test_run_pack_renders_slides_without_template(tmp_path: Path) -> None:
    pack_path = tmp_path / "pack.yaml"
    pack_path.write_text("", encoding="utf-8")

    pack = PackConfig(
        schema="test-pack",
        slides=[
            PackSlide(title="With Template", id="with-template", template="full_page_image", visual=PackVisualRef(ref="one.yaml")),
            PackSlide(title="No Template", id="no-template", visual=PackVisualRef(ref="two.yaml")),
        ],
    )

    visuals: Dict[str, BaseVisualConfig] = {
        "one.yaml": BaseVisualConfig(type="matrix"),
        "two.yaml": BaseVisualConfig(type="matrix"),
    }

    def _loader(path: Path, payload: Mapping[str, object] | None = None, stack: tuple[Path, ...] = ()) -> BaseVisualConfig:
        return visuals[path.name]

    pipeline = _StubPipeline()
    results = run_pack(
        pack_path,
        pack,
        project_root=pack_path.parent,
        output_root=tmp_path / "artefacts",
        base_options=PipelineOptions(),
        visual_loader=_loader,
        pipeline=cast(VisualPipeline[Any], pipeline),
        env=create_pack_jinja_env(),
    )

    png_paths = {result.png_path for result in results if result.png_path}
    assert (tmp_path / "artefacts" / "[01]_with_template.png") in png_paths
    assert (tmp_path / "artefacts" / "[02]_no_template.png") in png_paths


def test_run_pack_renders_placeholder_visuals(tmp_path: Path) -> None:
    pack_path = tmp_path / "pack.yaml"
    pack_path.write_text("", encoding="utf-8")

    pack = PackConfig(
        schema="test-pack",
        slides=[
            PackSlide(
                title="Two Up",
                id="two-up",
                template="two_up",
                placeholders={
                    "left_chart": PackPlaceholder(visual=PackVisualRef(ref="left.yaml")),
                    "right_chart": PackPlaceholder(visual=PackVisualRef(ref="right.yaml")),
                },
            )
        ],
    )

    visuals: Dict[str, BaseVisualConfig] = {
        "left.yaml": BaseVisualConfig(type="matrix"),
        "right.yaml": BaseVisualConfig(type="matrix"),
    }

    def _loader(path: Path, payload: Mapping[str, object] | None = None, stack: tuple[Path, ...] = ()) -> BaseVisualConfig:
        return visuals[path.name]

    pipeline = _StubPipeline()
    results = run_pack(
        pack_path,
        pack,
        project_root=pack_path.parent,
        output_root=tmp_path / "artefacts",
        base_options=PipelineOptions(),
        visual_loader=_loader,
        pipeline=cast(VisualPipeline[Any], pipeline),
        env=create_pack_jinja_env(),
    )

    png_paths = {result.png_path for result in results if result.png_path}
    expected_left = tmp_path / "artefacts" / "[01]_two_up__left_chart.png"
    expected_right = tmp_path / "artefacts" / "[01]_two_up__right_chart.png"
    assert expected_left in png_paths
    assert expected_right in png_paths


def test_run_pack_invokes_pptx_when_result_file(monkeypatch, tmp_path: Path) -> None:
    pack_path = tmp_path / "pack.yaml"
    pack_path.write_text("", encoding="utf-8")

    # Place a template where the resolver will find it.
    template_path = pack_path.parent / "pack_template.pptx"
    template_path.write_bytes(b"template")

    pack = PackConfig(
        schema="test-pack",
        slides=[
            PackSlide(title="Deck Slide", id="deck-slide", template="full_page_image", visual=PackVisualRef(ref="one.yaml")),
        ],
    )

    visuals: Dict[str, BaseVisualConfig] = {
        "one.yaml": BaseVisualConfig(type="matrix"),
    }

    def _loader(path: Path, payload: Mapping[str, object] | None = None, stack: tuple[Path, ...] = ()) -> BaseVisualConfig:
        return visuals[path.name]

    captured: Dict[str, object] = {}

    def _fake_assemble(**kwargs) -> None:
        captured.update(kwargs)

    monkeypatch.setattr("praeparo.pack.runner.assemble_pack_pptx", _fake_assemble)

    pipeline = _StubPipeline()
    base_options = PipelineOptions(metadata={"result_file": tmp_path / "deck.pptx"})
    run_pack(
        pack_path,
        pack,
        project_root=pack_path.parent,
        output_root=tmp_path / "artefacts",
        base_options=base_options,
        visual_loader=_loader,
        pipeline=cast(VisualPipeline[Any], pipeline),
        env=create_pack_jinja_env(),
    )

    assert captured, "assemble_pack_pptx should be invoked when result_file is set"
    assert captured["slide_pngs"]
    assert captured["placeholder_pngs"] == {}
    assert captured["result_path"] == tmp_path / "deck.pptx"
    assert captured["template_path"] == template_path


def test_run_pack_builds_pptx_with_template_only_slide(tmp_path: Path) -> None:
    pack_path = tmp_path / "pack.yaml"
    pack_path.write_text("", encoding="utf-8")

    template_path = tmp_path / "pack_template.pptx"
    prs = Presentation()
    picture_with_caption = prs.slide_layouts[8]

    # Single-image template uses a picture placeholder so the single-visual shorthand works.
    single = prs.slides.add_slide(picture_with_caption)
    image_ph = next(ph for ph in single.placeholders if ph.placeholder_format.type == PP_PLACEHOLDER.PICTURE)
    image_ph.name = "image"
    single.notes_slide.notes_text_frame.text = "TEMPLATE_TAG=single_image"

    # Two-up template relies on explicit pictures so placeholder binding matches by name.
    blank = prs.slide_layouts[6]
    two_up = prs.slides.add_slide(blank)
    tmp_img = tmp_path / "tmp.png"
    _write_png(tmp_img)
    left = two_up.shapes.add_picture(str(tmp_img), Inches(0.5), Inches(1), width=Inches(3), height=Inches(2.5))
    left.name = "left_chart"
    right = two_up.shapes.add_picture(str(tmp_img), Inches(4), Inches(1), width=Inches(3), height=Inches(2.5))
    right.name = "right_chart"
    two_up.notes_slide.notes_text_frame.text = "TEMPLATE_TAG=two_up"
    prs.save(template_path)
    tmp_img.unlink(missing_ok=True)
    result_path = tmp_path / "deck" / "governance.pptx"

    pack = PackConfig(
        schema="test-pack",
        slides=[
            PackSlide(title="Single Visual", id="single", template="single_image", visual=PackVisualRef(ref="one.yaml")),
            PackSlide(
                title="Placeholders",
                id="placeholders",
                template="two_up",
                placeholders={
                    "left_chart": PackPlaceholder(visual=PackVisualRef(ref="left.yaml")),
                    "right_chart": PackPlaceholder(visual=PackVisualRef(ref="right.yaml")),
                },
            ),
            PackSlide(title="Static Home", id="home", template="single_image"),
        ],
    )

    visuals: Dict[str, BaseVisualConfig] = {
        "one.yaml": BaseVisualConfig(type="matrix"),
        "left.yaml": BaseVisualConfig(type="matrix"),
        "right.yaml": BaseVisualConfig(type="matrix"),
    }

    def _loader(path: Path, payload: Mapping[str, object] | None = None, stack: tuple[Path, ...] = ()) -> BaseVisualConfig:
        return visuals[path.name]

    pipeline = _PngPipeline()
    base_options = PipelineOptions(metadata={"result_file": result_path, "pptx_template": template_path})
    run_pack(
        pack_path,
        pack,
        project_root=pack_path.parent,
        output_root=tmp_path / "artefacts",
        base_options=base_options,
        visual_loader=_loader,
        pipeline=cast(VisualPipeline[Any], pipeline),
        env=create_pack_jinja_env(),
    )

    deck = Presentation(result_path)
    assert len(deck.slides) == 3
    template_only_slide = deck.slides[2]
    picture_placeholders = []
    for shape in template_only_slide.shapes:
        if not getattr(shape, "is_placeholder", False):
            continue
        try:
            if shape.placeholder_format.type == PP_PLACEHOLDER.PICTURE:
                picture_placeholders.append(shape)
        except ValueError:
            continue
    assert picture_placeholders


def test_run_pack_renders_named_text_placeholders(tmp_path: Path) -> None:
    pack_path = tmp_path / "pack.yaml"
    pack_path.write_text("", encoding="utf-8")

    template_path = tmp_path / "pack_template.pptx"
    _build_two_up_template_with_text(template_path)

    result_path = tmp_path / "deck" / "text_placeholders.pptx"

    pack = PackConfig(
        schema="test-pack",
        context=PackContext.model_validate({"month": "2025-10-01"}),
        slides=[
            PackSlide(
                title="Two Up With Text",
                id="two-up-text",
                template="two_up",
                placeholders={
                    "display_date_text": PackPlaceholder(text="{{ month }}"),
                    "left_chart": PackPlaceholder(visual=PackVisualRef(ref="left.yaml")),
                    "right_chart": PackPlaceholder(visual=PackVisualRef(ref="right.yaml")),
                },
            )
        ],
    )

    visuals: Dict[str, BaseVisualConfig] = {
        "left.yaml": BaseVisualConfig(type="matrix"),
        "right.yaml": BaseVisualConfig(type="matrix"),
    }

    def _loader(path: Path, payload: Mapping[str, object] | None = None, stack: tuple[Path, ...] = ()) -> BaseVisualConfig:
        return visuals[path.name]

    pipeline = _PngPipeline()
    base_options = PipelineOptions(metadata={"result_file": result_path, "pptx_template": template_path})
    run_pack(
        pack_path,
        pack,
        project_root=pack_path.parent,
        output_root=tmp_path / "artefacts",
        base_options=base_options,
        visual_loader=_loader,
        pipeline=cast(VisualPipeline[Any], pipeline),
        env=create_pack_jinja_env(),
    )

    assert len(pipeline.calls) == 2

    deck = Presentation(result_path)
    assert len(deck.slides) == 1
    slide = deck.slides[0]
    text_shapes = [shape for shape in slide.shapes if getattr(shape, "name", None) == "display_date_text"]
    assert text_shapes
    assert text_shapes[0].text_frame.text.strip() == "2025-10-01"


def test_run_pack_binds_static_slide_image(tmp_path: Path) -> None:
    pack_path = tmp_path / "pack.yaml"
    pack_path.write_text("", encoding="utf-8")

    template_path = tmp_path / "pack_template.pptx"
    prs = Presentation()
    picture_with_caption = prs.slide_layouts[8]

    single = prs.slides.add_slide(picture_with_caption)
    image_ph = next(ph for ph in single.placeholders if ph.placeholder_format.type == PP_PLACEHOLDER.PICTURE)
    image_ph.name = "image"
    single.notes_slide.notes_text_frame.text = "TEMPLATE_TAG=single_image"

    blank = prs.slide_layouts[6]
    two_up = prs.slides.add_slide(blank)
    tmp_img = tmp_path / "tmp.png"
    _write_png(tmp_img)
    left = two_up.shapes.add_picture(str(tmp_img), Inches(0.5), Inches(1), width=Inches(3), height=Inches(2.5))
    left.name = "left_chart"
    right = two_up.shapes.add_picture(str(tmp_img), Inches(4), Inches(1), width=Inches(3), height=Inches(2.5))
    right.name = "right_chart"
    two_up.notes_slide.notes_text_frame.text = "TEMPLATE_TAG=two_up"

    prs.save(template_path)
    tmp_img.unlink(missing_ok=True)
    result_path = tmp_path / "deck" / "static_home.pptx"

    logo_path = pack_path.parent / "logo.png"
    _write_coloured_png(logo_path, colour=(0, 255, 0, 255))

    pack = PackConfig(
        schema="test-pack",
        slides=[
            PackSlide(
                title="Static Home",
                id="home",
                template="single_image",
                image="logo.png",
            )
        ],
    )

    base_options = PipelineOptions(metadata={"result_file": result_path, "pptx_template": template_path})
    run_pack(
        pack_path,
        pack,
        project_root=pack_path.parent,
        output_root=tmp_path / "artefacts",
        base_options=base_options,
        pipeline=cast(VisualPipeline[Any], _PngPipeline()),
        env=create_pack_jinja_env(),
    )

    deck = Presentation(result_path)
    assert len(deck.slides) == 1
    pictures = _picture_blobs_by_name(deck.slides[0])
    assert pictures.get("image") == logo_path.read_bytes()


def test_run_pack_binds_templated_static_slide_image(tmp_path: Path) -> None:
    pack_path = tmp_path / "pack.yaml"
    pack_path.write_text("", encoding="utf-8")

    template_path = tmp_path / "pack_template.pptx"
    _build_pack_template(template_path)
    result_path = tmp_path / "deck" / "templated_static_home.pptx"

    logo_path = pack_path.parent / "logo.png"
    _write_coloured_png(logo_path, colour=(0, 255, 0, 255))

    pack = PackConfig(
        schema="test-pack",
        context=PackContext.model_validate({"logo": "logo.png"}),
        slides=[
            PackSlide(
                title="Static Home",
                id="home",
                template="single_image",
                image="{{ logo }}",
            )
        ],
    )

    base_options = PipelineOptions(metadata={"result_file": result_path, "pptx_template": template_path})
    run_pack(
        pack_path,
        pack,
        project_root=pack_path.parent,
        output_root=tmp_path / "artefacts",
        base_options=base_options,
        pipeline=cast(VisualPipeline[Any], _PngPipeline()),
        env=create_pack_jinja_env(),
    )

    deck = Presentation(result_path)
    assert len(deck.slides) == 1
    pictures = _picture_blobs_by_name(deck.slides[0])
    assert pictures.get("image") == logo_path.read_bytes()


def test_run_pack_binds_static_placeholder_image(tmp_path: Path) -> None:
    pack_path = tmp_path / "pack.yaml"
    pack_path.write_text("", encoding="utf-8")

    template_path = tmp_path / "pack_template.pptx"
    result_path = tmp_path / "deck" / "static_placeholder.pptx"

    _build_pack_template(template_path)

    logo_path = pack_path.parent / "logo.png"
    _write_coloured_png(logo_path, colour=(0, 0, 255, 255))

    pack = PackConfig(
        schema="test-pack",
        slides=[
            PackSlide(
                title="Two Up",
                id="two-up",
                template="two_up",
                placeholders={
                    "left_chart": PackPlaceholder(visual=PackVisualRef(ref="left.yaml")),
                    "right_chart": PackPlaceholder(image="logo.png"),
                },
            )
        ],
    )

    visuals: Dict[str, BaseVisualConfig] = {"left.yaml": BaseVisualConfig(type="matrix")}

    def _loader(path: Path, payload: Mapping[str, object] | None = None, stack: tuple[Path, ...] = ()) -> BaseVisualConfig:
        return visuals[path.name]

    pipeline = _PngPipeline()
    base_options = PipelineOptions(metadata={"result_file": result_path, "pptx_template": template_path})
    run_pack(
        pack_path,
        pack,
        project_root=pack_path.parent,
        output_root=tmp_path / "artefacts",
        base_options=base_options,
        visual_loader=_loader,
        pipeline=cast(VisualPipeline[Any], pipeline),
        env=create_pack_jinja_env(),
    )

    deck = Presentation(result_path)
    assert len(deck.slides) == 1
    pictures = _picture_blobs_by_name(deck.slides[0])
    left_png = (tmp_path / "artefacts" / f"[01]_{slugify('two-up')}__{slugify('left_chart')}.png").read_bytes()
    assert pictures.get("left_chart") == left_png
    assert pictures.get("right_chart") == logo_path.read_bytes()


def test_run_pack_binds_templated_static_placeholder_image(tmp_path: Path) -> None:
    pack_path = tmp_path / "pack.yaml"
    pack_path.write_text("", encoding="utf-8")

    template_path = tmp_path / "pack_template.pptx"
    result_path = tmp_path / "deck" / "templated_static_placeholder.pptx"
    _build_pack_template(template_path)

    logo_path = pack_path.parent / "logo.png"
    _write_coloured_png(logo_path, colour=(0, 0, 255, 255))

    pack = PackConfig(
        schema="test-pack",
        context=PackContext.model_validate({"logo": "logo.png"}),
        slides=[
            PackSlide(
                title="Two Up",
                id="two-up",
                template="two_up",
                placeholders={
                    "left_chart": PackPlaceholder(visual=PackVisualRef(ref="left.yaml")),
                    "right_chart": PackPlaceholder(image="{{ logo }}"),
                },
            )
        ],
    )

    visuals: Dict[str, BaseVisualConfig] = {"left.yaml": BaseVisualConfig(type="matrix")}

    def _loader(path: Path, payload: Mapping[str, object] | None = None, stack: tuple[Path, ...] = ()) -> BaseVisualConfig:
        return visuals[path.name]

    pipeline = _PngPipeline()
    base_options = PipelineOptions(metadata={"result_file": result_path, "pptx_template": template_path})
    run_pack(
        pack_path,
        pack,
        project_root=pack_path.parent,
        output_root=tmp_path / "artefacts",
        base_options=base_options,
        visual_loader=_loader,
        pipeline=cast(VisualPipeline[Any], pipeline),
        env=create_pack_jinja_env(),
    )

    deck = Presentation(result_path)
    assert len(deck.slides) == 1
    pictures = _picture_blobs_by_name(deck.slides[0])
    left_png = (tmp_path / "artefacts" / f"[01]_{slugify('two-up')}__{slugify('left_chart')}.png").read_bytes()
    assert pictures.get("left_chart") == left_png
    assert pictures.get("right_chart") == logo_path.read_bytes()


def test_run_pack_errors_on_missing_static_images(tmp_path: Path) -> None:
    pack_path = tmp_path / "pack.yaml"
    pack_path.write_text("", encoding="utf-8")

    slide_only_pack = PackConfig(
        schema="test-pack",
        slides=[
            PackSlide(
                title="Broken Slide",
                id="broken",
                template="single_image",
                image="missing.png",
            ),
        ],
    )

    with pytest.raises(ValueError, match="Static slide image not found"):
        run_pack(
            pack_path,
            slide_only_pack,
            project_root=pack_path.parent,
            output_root=tmp_path / "artefacts",
            base_options=PipelineOptions(),
            pipeline=cast(VisualPipeline[Any], _PngPipeline()),
            env=create_pack_jinja_env(),
        )

    placeholder_pack = PackConfig(
        schema="test-pack",
        slides=[
            PackSlide(
                title="Broken Placeholder",
                id="broken-placeholder",
                template="two_up",
                placeholders={
                    "logo": PackPlaceholder(image="missing_logo.png"),
                },
            )
        ],
    )

    with pytest.raises(ValueError, match="Static placeholder image not found"):
        run_pack(
            pack_path,
            placeholder_pack,
            project_root=pack_path.parent,
            output_root=tmp_path / "artefacts",
            base_options=PipelineOptions(),
            pipeline=cast(VisualPipeline[Any], _PngPipeline()),
            env=create_pack_jinja_env(),
        )


def test_run_pack_errors_when_templated_static_image_resolves_empty(tmp_path: Path) -> None:
    pack_path = tmp_path / "pack.yaml"
    pack_path.write_text("", encoding="utf-8")

    template_path = tmp_path / "pack_template.pptx"
    _build_pack_template(template_path)

    pack = PackConfig(
        schema="test-pack",
        slides=[
            PackSlide(
                title="Broken Templated Image",
                id="broken-templated-image",
                template="single_image",
                image="{{ logo }}",
            ),
        ],
    )

    with pytest.raises(ValueError, match="resolved to an empty path"):
        run_pack(
            pack_path,
            pack,
            project_root=pack_path.parent,
            output_root=tmp_path / "artefacts",
            base_options=PipelineOptions(),
            pipeline=cast(VisualPipeline[Any], _PngPipeline()),
            env=create_pack_jinja_env(),
        )


def test_pack_models_validate_static_image_rules() -> None:
    with pytest.raises(ValidationError):
        PackPlaceholder(visual=PackVisualRef(ref="one.yaml"), image="logo.png")

    with pytest.raises(ValidationError):
        PackPlaceholder(visual=PackVisualRef(ref="one.yaml"), text="hello")

    with pytest.raises(ValidationError):
        PackPlaceholder()

    with pytest.raises(ValidationError):
        PackSlide(title="Image Missing Template", image="logo.png")

    with pytest.raises(ValidationError):
        PackSlide(title="Image With Visual", template="single", visual=PackVisualRef(ref="one.yaml"), image="logo.png")


def test_pack_slide_placeholders_support_string_shorthand() -> None:
    pack = PackConfig.model_validate(
        {
            "schema": "test-pack",
            "slides": [
                {
                    "title": "Home",
                    "template": "home",
                    "placeholders": {
                        "logo": "./assets/amp_logo.png",
                        "text": "{{ display_date }}",
                    },
                }
            ],
        }
    )

    placeholders = pack.slides[0].placeholders
    assert placeholders is not None
    assert placeholders["logo"].image == "./assets/amp_logo.png"
    assert placeholders["logo"].text is None
    assert placeholders["text"].text == "{{ display_date }}"
    assert placeholders["text"].image is None


def test_pack_visual_ref_requires_ref_or_type() -> None:
    with pytest.raises(ValidationError):
        PackVisualRef()

    with pytest.raises(ValidationError):
        PackVisualRef(ref="one.yaml", type="two.yaml")


def test_pack_visual_ref_series_operations_require_ref() -> None:
    with pytest.raises(ValidationError):
        PackVisualRef.model_validate({"type": "matrix", "series_add": [{"id": "other_lender"}]})


def test_pack_visual_ref_rejects_series_and_series_operations_mix() -> None:
    with pytest.raises(ValidationError):
        PackVisualRef.model_validate(
            {
                "ref": "one.yaml",
                "series": [{"id": "customer_lender"}],
                "series_add": [{"id": "other_lender"}],
            }
        )


def test_run_pack_supports_python_visual_ref(tmp_path: Path, monkeypatch) -> None:
    fixture_path = Path(__file__).parent.parent / "fixtures" / "python_visuals" / "pack_python_visual.py"
    python_visual_path = tmp_path / "visuals" / "pack_python_visual.py"
    python_visual_path.parent.mkdir(parents=True, exist_ok=True)
    python_visual_path.write_text(fixture_path.read_text(encoding="utf-8"), encoding="utf-8")

    pack_path = tmp_path / "pack.yaml"
    pack_path.write_text("", encoding="utf-8")

    pack = PackConfig(
        schema="test-pack",
        context=PackContext.model_validate({"title": "Demo"}),
        slides=[
            PackSlide(
                title="Python Visual Slide",
                template="full_page_image",
                visual=PackVisualRef(ref=str(python_visual_path)),
            )
        ],
    )

    captured: dict[str, object] = {}

    def fake_write_image(self, output_path, *, scale=2.0, **_: object) -> None:
        Path(output_path).parent.mkdir(parents=True, exist_ok=True)
        Path(output_path).write_bytes(b"PNG")
        captured["scale"] = scale

    monkeypatch.setattr(go.Figure, "write_image", fake_write_image, raising=False)

    output_root = tmp_path / "artefacts"
    pipeline = VisualPipeline(planner_provider=build_default_query_planner_provider())

    results = run_pack(
        pack_path,
        pack,
        project_root=pack_path.parent,
        output_root=output_root,
        base_options=PipelineOptions(),
        pipeline=pipeline,
        env=create_pack_jinja_env(),
    )

    assert results
    expected_png = output_root / "[01]_python_visual_slide.png"
    assert expected_png.exists()
    assert results[0].png_path == expected_png
    assert captured.get("scale") == 2.0


def test_default_metrics_root_prefers_registry_metrics(tmp_path: Path) -> None:
    metrics_dir = tmp_path / "repo" / "registry" / "metrics"
    metrics_dir.mkdir(parents=True)
    pack_path = metrics_dir.parent.parent / "customers" / "amp" / "pack.yaml"
    pack_path.parent.mkdir(parents=True)
    pack_path.write_text("{}", encoding="utf-8")

    from praeparo.datasets.context import resolve_default_metrics_root_for_pack

    resolved = resolve_default_metrics_root_for_pack(pack_path)
    assert resolved == metrics_dir.resolve()


def test_run_pack_supports_inline_visual_config(tmp_path: Path, monkeypatch) -> None:
    fixture_path = Path(__file__).parent.parent / "fixtures" / "python_visuals" / "pack_python_visual.py"
    python_visual_path = tmp_path / "visuals" / "pack_python_visual.py"
    python_visual_path.parent.mkdir(parents=True, exist_ok=True)
    python_visual_path.write_text(fixture_path.read_text(encoding="utf-8"), encoding="utf-8")

    pack_path = tmp_path / "pack.yaml"
    pack_path.write_text("", encoding="utf-8")

    pack = PackConfig(
        schema="test-pack",
        context=PackContext.model_validate({"title": "Inline Demo"}),
        slides=[
            PackSlide(
                title="Inline Visual Slide",
                template="full_page_image",
                visual=PackVisualRef(type=str(python_visual_path)),
            )
        ],
    )

    captured: dict[str, object] = {}

    def fake_write_image(self, output_path, *, scale=2.0, **_: object) -> None:
        Path(output_path).parent.mkdir(parents=True, exist_ok=True)
        Path(output_path).write_bytes(b"PNG")
        captured["path"] = output_path

    monkeypatch.setattr(go.Figure, "write_image", fake_write_image, raising=False)

    output_root = tmp_path / "artefacts"
    pipeline = VisualPipeline(planner_provider=build_default_query_planner_provider())

    results = run_pack(
        pack_path,
        pack,
        project_root=pack_path.parent,
        output_root=output_root,
        base_options=PipelineOptions(),
        pipeline=pipeline,
        env=create_pack_jinja_env(),
    )

    assert results
    expected_png = output_root / "[01]_inline_visual_slide.png"
    assert expected_png.exists()
    assert results[0].png_path == expected_png
    assert captured.get("path") == expected_png


def test_restitch_pack_pptx_honours_templated_titles(tmp_path: Path) -> None:
    pack_path = tmp_path / "pack.yaml"
    pack_path.write_text("", encoding="utf-8")

    template_path = tmp_path / "pack_template.pptx"
    prs = Presentation()
    picture_with_caption = prs.slide_layouts[8]
    single = prs.slides.add_slide(picture_with_caption)
    image_ph = next(ph for ph in single.placeholders if ph.placeholder_format.type == PP_PLACEHOLDER.PICTURE)
    image_ph.name = "image"
    single.notes_slide.notes_text_frame.text = "TEMPLATE_TAG=single_image"
    prs.save(template_path)

    result_path = tmp_path / "deck" / "restitched.pptx"

    pack = PackConfig(
        schema="test-pack",
        context=PackContext.model_validate({"customer": "AMP"}),
        slides=[
            PackSlide(
                title="{{customer}} Dashboard",
                template="single_image",
                visual=PackVisualRef(ref="visual.yaml"),
            ),
        ],
    )

    output_root = tmp_path / "artefacts"
    output_root.mkdir(parents=True, exist_ok=True)

    slug = slugify("AMP Dashboard")
    png_path = output_root / f"[01]_{slug}.png"
    _write_coloured_png(png_path, colour=(0, 0, 255, 255))

    base_options = PipelineOptions(metadata={"pptx_template": template_path, "result_file": result_path})
    restitch_pack_pptx(
        pack_path,
        pack,
        output_root=output_root,
        result_file=result_path,
        base_options=base_options,
    )

    prs = Presentation(result_path)
    assert len(prs.slides) == 1
    blobs = _picture_blobs_by_name(prs.slides[0])
    assert blobs.get("image") == png_path.read_bytes()


def test_restitch_pack_pptx_binds_templated_static_placeholder_image(tmp_path: Path) -> None:
    pack_path = tmp_path / "pack.yaml"
    pack_path.write_text("", encoding="utf-8")

    template_path = tmp_path / "pack_template.pptx"
    _build_pack_template(template_path)

    result_path = tmp_path / "deck" / "restitched_templated_static_placeholder.pptx"

    logo_path = pack_path.parent / "logo.png"
    _write_coloured_png(logo_path, colour=(0, 255, 0, 255))

    pack = PackConfig(
        schema="test-pack",
        context=PackContext.model_validate({"logo": "logo.png"}),
        slides=[
            PackSlide(
                title="Two Up Slide",
                template="two_up",
                placeholders={
                    "right_chart": PackPlaceholder(image="{{ logo }}"),
                },
            ),
        ],
    )

    output_root = tmp_path / "artefacts"
    output_root.mkdir(parents=True, exist_ok=True)

    base_options = PipelineOptions(metadata={"pptx_template": template_path, "result_file": result_path})
    restitch_pack_pptx(
        pack_path,
        pack,
        output_root=output_root,
        result_file=result_path,
        base_options=base_options,
    )

    prs = Presentation(result_path)
    assert len(prs.slides) == 1
    blobs = _picture_blobs_by_name(prs.slides[0])
    assert blobs.get("right_chart") == logo_path.read_bytes()


def test_restitch_pack_pptx_reuses_existing_assets(tmp_path: Path) -> None:
    pack_path = tmp_path / "pack.yaml"
    pack_path.write_text("", encoding="utf-8")

    template_path = tmp_path / "pack_template.pptx"
    prs = Presentation()
    picture_with_caption = prs.slide_layouts[8]

    single = prs.slides.add_slide(picture_with_caption)
    image_ph = next(ph for ph in single.placeholders if ph.placeholder_format.type == PP_PLACEHOLDER.PICTURE)
    image_ph.name = "image"
    single.notes_slide.notes_text_frame.text = "TEMPLATE_TAG=single_image"

    blank = prs.slide_layouts[6]
    two_up = prs.slides.add_slide(blank)
    tmp_img = tmp_path / "tmp.png"
    _write_png(tmp_img)
    left = two_up.shapes.add_picture(str(tmp_img), Inches(0.5), Inches(1), width=Inches(3), height=Inches(2.5))
    left.name = "left_chart"
    right = two_up.shapes.add_picture(str(tmp_img), Inches(4), Inches(1), width=Inches(3), height=Inches(2.5))
    right.name = "right_chart"
    two_up.notes_slide.notes_text_frame.text = "TEMPLATE_TAG=two_up"

    prs.save(template_path)
    tmp_img.unlink(missing_ok=True)
    result_path = tmp_path / "deck" / "restitched.pptx"

    pack = PackConfig(
        schema="test-pack",
        slides=[
            PackSlide(
                title="Static Slide",
                template="single_image",
                image="static.png",
            ),
            PackSlide(
                title="Two Up Slide",
                template="two_up",
                placeholders={
                    "left_chart": PackPlaceholder(visual=PackVisualRef(ref="left.yaml")),
                    "right_chart": PackPlaceholder(image="right.png"),
                },
            ),
            PackSlide(
                title="Visual Slide",
                template="single_image",
                visual=PackVisualRef(ref="visual.yaml"),
            ),
        ],
    )

    pack_root = pack_path.parent
    pack_root.mkdir(parents=True, exist_ok=True)
    static_path = pack_root / "static.png"
    _write_png(static_path)
    right_placeholder_path = pack_root / "right.png"
    _write_coloured_png(right_placeholder_path, colour=(0, 255, 0, 255))

    output_root = tmp_path / "artefacts"
    output_root.mkdir(parents=True, exist_ok=True)

    static_slug = slugify("Static Slide")
    two_up_slug = slugify("Two Up Slide")
    visual_slug = slugify("Visual Slide")

    _write_coloured_png(output_root / f"[01]_{static_slug}.png", colour=(0, 0, 255, 255))
    _write_coloured_png(output_root / f"[02]_{two_up_slug}__left_chart.png", colour=(255, 0, 0, 255))
    _write_coloured_png(output_root / f"[03]_{visual_slug}.png", colour=(128, 128, 128, 255))

    base_options = PipelineOptions(metadata={"pptx_template": template_path, "result_file": result_path})
    restitch_pack_pptx(
        pack_path,
        pack,
        output_root=output_root,
        result_file=result_path,
        base_options=base_options,
    )

    prs = Presentation(result_path)
    assert len(prs.slides) == 3

    static_blobs = _picture_blobs_by_name(prs.slides[0])
    assert static_blobs.get("image") == static_path.read_bytes()

    two_up_blobs = _picture_blobs_by_name(prs.slides[1])
    left_png = (output_root / f"[02]_{two_up_slug}__left_chart.png").read_bytes()
    assert two_up_blobs.get("left_chart") == left_png
    assert two_up_blobs.get("right_chart") == right_placeholder_path.read_bytes()

    visual_blobs = _picture_blobs_by_name(prs.slides[2])
    assert visual_blobs.get("image") == (output_root / f"[03]_{visual_slug}.png").read_bytes()
