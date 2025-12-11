from __future__ import annotations

from pathlib import Path

import base64

from pptx import Presentation
from pptx.enum.shapes import PP_PLACEHOLDER
from pptx.util import Inches

from praeparo.pack.pptx import PlaceholderSize, resolve_template_geometry


def _write_png(path: Path) -> None:
    data = base64.b64decode(
        "iVBORw0KGgoAAAANSUhEUgAAAAEAAAABCAYAAAAfFcSJAAAADUlEQVR42mP8z8BQDwAEAAL/AXx5lXcAAAAASUVORK5CYII="
    )
    path.write_bytes(data)


def _build_template(path: Path) -> None:
    prs = Presentation()
    blank = prs.slide_layouts[6]
    picture_with_caption = prs.slide_layouts[8]

    tmp = path.parent / "tmp.png"
    _write_png(tmp)

    single = prs.slides.add_slide(picture_with_caption)
    image_ph = next(ph for ph in single.placeholders if ph.placeholder_format.type == PP_PLACEHOLDER.PICTURE)
    image_ph.name = "image"
    single.notes_slide.notes_text_frame.text = "TEMPLATE_TAG=single_image"

    two_up = prs.slides.add_slide(blank)
    left = two_up.shapes.add_picture(str(tmp), Inches(0.5), Inches(1), width=Inches(3), height=Inches(2.5))
    left.name = "left_chart"
    right = two_up.shapes.add_picture(str(tmp), Inches(4), Inches(1), width=Inches(3), height=Inches(2.5))
    right.name = "right_chart"
    two_up.notes_slide.notes_text_frame.text = "TEMPLATE_TAG=two_up"

    prs.save(path)
    tmp.unlink(missing_ok=True)


def test_resolve_template_geometry(tmp_path: Path) -> None:
    template = tmp_path / "pack_template.pptx"
    _build_template(template)

    slide_sizes, placeholder_sizes = resolve_template_geometry(template)

    assert slide_sizes["single_image"].width_px > 0
    assert slide_sizes["single_image"].height_px > 0

    left = placeholder_sizes[("two_up", "left_chart")]
    right = placeholder_sizes[("two_up", "right_chart")]

    assert left.width_px > 0 and left.height_px > 0
    assert right.width_px > 0 and right.height_px > 0
    assert isinstance(left, PlaceholderSize)
    assert isinstance(right, PlaceholderSize)
    assert (left.width_px, left.height_px) == (right.width_px, right.height_px)
