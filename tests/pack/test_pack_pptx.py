from __future__ import annotations

import struct
import zlib
import zipfile
from pathlib import Path
import xml.etree.ElementTree as ET

import pptx.enum.shapes as pptx_shapes
from pptx.dml.color import RGBColor
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.util import Inches

from praeparo.pack.pptx import assemble_pack_pptx
from praeparo.models import PackConfig, PackPlaceholder, PackSlide, PackVisualRef
from praeparo.visuals.dax.planner_core import slugify

PP_PLACEHOLDER_ENUM = getattr(
    pptx_shapes,
    "PP_PLACEHOLDER_TYPE",
    getattr(pptx_shapes, "PP_PLACEHOLDER"),
)
MSO_PICTURE = getattr(MSO_SHAPE_TYPE, "PICTURE")


def _write_png(path: Path, width: int = 1, height: int = 1) -> None:
    """Write a solid-colour PNG of the requested dimensions (no external deps)."""

    raw_rows = [b"\x00" + b"\x00" * (width * 3) for _ in range(height)]
    raw = b"".join(raw_rows)
    compressed = zlib.compress(raw)

    def chunk(tag: bytes, payload: bytes) -> bytes:
        return struct.pack(">I", len(payload)) + tag + payload + struct.pack(">I", zlib.crc32(tag + payload) & 0xFFFFFFFF)

    signature = b"\x89PNG\r\n\x1a\n"
    ihdr = chunk(b"IHDR", struct.pack(">IIBBBBB", width, height, 8, 2, 0, 0, 0))
    idat = chunk(b"IDAT", compressed)
    iend = chunk(b"IEND", b"")

    path.write_bytes(signature + ihdr + idat + iend)


def _build_template(path: Path) -> None:
    prs = Presentation()
    picture_with_caption = prs.slide_layouts[8]
    tmp_img = path.parent / "tmp.png"
    _write_png(tmp_img)

    # Single-image template uses the built-in picture placeholder layout so we exercise placeholder logic.
    single = prs.slides.add_slide(picture_with_caption)
    picture_placeholder = next(ph for ph in single.placeholders if ph.placeholder_format.type == PP_PLACEHOLDER_ENUM.PICTURE)
    picture_placeholder.name = "image"
    picture_placeholder.left = Inches(1)
    picture_placeholder.top = Inches(1)
    picture_placeholder.width = Inches(4)
    picture_placeholder.height = Inches(3)
    single.notes_slide.notes_text_frame.text = "TEMPLATE_TAG=single_image"

    # Two-up template
    blank = prs.slide_layouts[6]
    two_up = prs.slides.add_slide(blank)
    left = two_up.shapes.add_picture(str(tmp_img), Inches(0.5), Inches(1), width=Inches(3), height=Inches(2.5))
    left.name = "left_chart"
    right = two_up.shapes.add_picture(str(tmp_img), Inches(4), Inches(1), width=Inches(3), height=Inches(2.5))
    right.name = "right_chart"
    two_up.notes_slide.notes_text_frame.text = "TEMPLATE_TAG=two_up"

    prs.save(path)
    tmp_img.unlink(missing_ok=True)


def _build_template_with_background(path: Path) -> None:
    prs = Presentation()
    picture_with_caption = prs.slide_layouts[8]

    slide = prs.slides.add_slide(picture_with_caption)
    tmp_img = path.parent / "tmp_background.png"
    _write_png(tmp_img, width=10, height=10)
    background = slide.shapes.add_picture(str(tmp_img), Inches(0), Inches(0), width=prs.slide_width, height=prs.slide_height)
    background.name = "background"

    logo_placeholder = next(ph for ph in slide.placeholders if ph.placeholder_format.type == PP_PLACEHOLDER_ENUM.PICTURE)
    logo_placeholder.name = "logo"
    logo_placeholder.left = Inches(1)
    logo_placeholder.top = Inches(1)
    logo_placeholder.width = Inches(2)
    logo_placeholder.height = Inches(1.5)

    slide.notes_slide.notes_text_frame.text = "TEMPLATE_TAG=home"
    prs.save(path)
    tmp_img.unlink(missing_ok=True)


def _build_template_with_slide_background_fill(path: Path) -> None:
    """Create a template where the home slide background is a slide-level image fill."""

    prs = Presentation()
    picture_with_caption = prs.slide_layouts[8]
    slide = prs.slides.add_slide(picture_with_caption)

    tmp_img = path.parent / "tmp_background_fill.png"
    _write_png(tmp_img, width=10, height=10)
    background_shape = slide.shapes.add_picture(
        str(tmp_img),
        Inches(0),
        Inches(0),
        width=prs.slide_width,
        height=prs.slide_height,
    )
    background_shape.name = "background"

    logo_placeholder = next(ph for ph in slide.placeholders if ph.placeholder_format.type == PP_PLACEHOLDER_ENUM.PICTURE)
    logo_placeholder.name = "logo"
    logo_placeholder.left = Inches(1)
    logo_placeholder.top = Inches(1)
    logo_placeholder.width = Inches(2)
    logo_placeholder.height = Inches(1.5)

    slide.notes_slide.notes_text_frame.text = "TEMPLATE_TAG=home"
    prs.save(path)
    tmp_img.unlink(missing_ok=True)

    namespaces = {
        "p": "http://schemas.openxmlformats.org/presentationml/2006/main",
        "a": "http://schemas.openxmlformats.org/drawingml/2006/main",
        "r": "http://schemas.openxmlformats.org/officeDocument/2006/relationships",
    }

    with zipfile.ZipFile(path, "r") as archive:
        entries = {name: archive.read(name) for name in archive.namelist()}

    slide_name = "ppt/slides/slide1.xml"
    root = ET.fromstring(entries[slide_name])
    c_sld = root.find("p:cSld", namespaces)
    assert c_sld is not None

    sp_tree = c_sld.find("p:spTree", namespaces)
    assert sp_tree is not None

    background_pic = sp_tree.find("p:pic", namespaces)
    assert background_pic is not None

    blip = background_pic.find(".//a:blip", namespaces)
    assert blip is not None
    embed_attr = f"{{{namespaces['r']}}}embed"
    rel_id = blip.attrib.get(embed_attr)
    assert rel_id is not None

    sp_tree.remove(background_pic)

    existing_bg = c_sld.find("p:bg", namespaces)
    if existing_bg is not None:
        c_sld.remove(existing_bg)

    bg = ET.Element(f"{{{namespaces['p']}}}bg")
    bg_pr = ET.SubElement(bg, f"{{{namespaces['p']}}}bgPr")
    blip_fill = ET.SubElement(bg_pr, f"{{{namespaces['a']}}}blipFill")
    blip_bg = ET.SubElement(blip_fill, f"{{{namespaces['a']}}}blip")
    blip_bg.set(embed_attr, rel_id)
    ET.SubElement(blip_fill, f"{{{namespaces['a']}}}srcRect")
    stretch = ET.SubElement(blip_fill, f"{{{namespaces['a']}}}stretch")
    ET.SubElement(stretch, f"{{{namespaces['a']}}}fillRect")
    ET.SubElement(bg_pr, f"{{{namespaces['a']}}}effectLst")
    c_sld.insert(0, bg)

    entries[slide_name] = ET.tostring(root, encoding="utf-8", xml_declaration=True)

    with zipfile.ZipFile(path, "w", zipfile.ZIP_DEFLATED) as archive:
        for name, payload in entries.items():
            archive.writestr(name, payload)


def test_assemble_pptx_single_placeholder_shorthand(tmp_path: Path) -> None:
    template_path = tmp_path / "template.pptx"
    _build_template(template_path)

    template_layout_name = Presentation(template_path).slides[0].slide_layout.name

    pack = PackConfig(
        schema="test-pack",
        slides=[
            PackSlide(
                title="With Template",
                id="with_template",
                template="single_image",
                visual=PackVisualRef(ref="vis.yaml"),
            )
        ],
    )

    slide_slug = slugify("with_template")
    png_path = tmp_path / f"{slide_slug}.png"
    _write_png(png_path)

    out = tmp_path / "deck.pptx"
    assemble_pack_pptx(
        pack=pack,
        results=[],
        slide_pngs={slide_slug: png_path},
        placeholder_pngs={},
        result_path=out,
        template_path=template_path,
    )

    prs = Presentation(out)
    assert len(prs.slides) == 1
    picture_shapes = [s for s in prs.slides[0].shapes if s.shape_type == MSO_PICTURE]
    assert picture_shapes  # image placed
    assert prs.slides[0].slide_layout.name == template_layout_name


def test_replace_picture_fits_and_centers(tmp_path: Path) -> None:
    template_path = tmp_path / "template.pptx"
    _build_template(template_path)

    template_prs = Presentation(template_path)
    placeholder = next(shape for shape in template_prs.slides[0].shapes if shape.name == "image")
    box_left, box_top, box_width, box_height = placeholder.left, placeholder.top, placeholder.width, placeholder.height

    pack = PackConfig(
        schema="test-pack",
        slides=[
            PackSlide(
                title="With Template",
                id="with_template",
                template="single_image",
                visual=PackVisualRef(ref="vis.yaml"),
            )
        ],
    )

    slide_slug = slugify("with_template")
    png_path = tmp_path / f"{slide_slug}.png"
    _write_png(png_path, width=400, height=100)  # wide image to test fitting

    out = tmp_path / "deck.pptx"
    assemble_pack_pptx(
        pack=pack,
        results=[],
        slide_pngs={slide_slug: png_path},
        placeholder_pngs={},
        result_path=out,
        template_path=template_path,
    )

    prs = Presentation(out)
    picture_shapes = [s for s in prs.slides[0].shapes if s.shape_type == MSO_PICTURE]
    assert len(picture_shapes) == 1

    pic = picture_shapes[0]
    ratio = pic.width / pic.height
    original_ratio = 400 / 100
    assert abs(ratio / original_ratio - 1.0) < 0.01

    assert pic.width <= box_width
    assert pic.height <= box_height

    left_gap = pic.left - box_left
    top_gap = pic.top - box_top
    right_gap = box_left + box_width - (pic.left + pic.width)
    bottom_gap = box_top + box_height - (pic.top + pic.height)

    assert left_gap >= 0 and top_gap >= 0
    assert right_gap >= 0 and bottom_gap >= 0
    assert abs(left_gap - right_gap) <= 10
    assert abs(top_gap - bottom_gap) <= 10


def test_assemble_pptx_multi_placeholder_and_skip_no_template(tmp_path: Path) -> None:
    template_path = tmp_path / "template.pptx"
    _build_template(template_path)

    pack = PackConfig(
        schema="test-pack",
        slides=[
            PackSlide(
                title="Two Up",
                id="two_up",
                template="two_up",
                placeholders={
                    "left_chart": PackPlaceholder(visual=PackVisualRef(ref="left.yaml")),
                    "right_chart": PackPlaceholder(visual=PackVisualRef(ref="right.yaml")),
                },
            ),
            PackSlide(
                title="No Template",
                id="no_template",
                visual=PackVisualRef(ref="plain.yaml"),
            ),
            PackSlide(
                title="Another Template",
                id="another_template",
                template="single_image",
                visual=PackVisualRef(ref="third.yaml"),
            ),
        ],
    )

    two_slug = slugify("two_up")
    left_png = tmp_path / f"{two_slug}_left.png"
    right_png = tmp_path / f"{two_slug}_right.png"
    _write_png(left_png)
    _write_png(right_png)
    placeholder_pngs = {two_slug: {f"{two_slug}__{slugify('left_chart')}": left_png, f"{two_slug}__{slugify('right_chart')}": right_png}}

    third_slug = slugify("another_template")
    third_png = tmp_path / f"{third_slug}.png"
    _write_png(third_png)

    out = tmp_path / "deck_multi.pptx"
    assemble_pack_pptx(
        pack=pack,
        results=[],
        slide_pngs={third_slug: third_png},
        placeholder_pngs=placeholder_pngs,
        result_path=out,
        template_path=template_path,
    )

    prs = Presentation(out)
    assert len(prs.slides) == 2  # no-template slide skipped

    first_slide_pictures = [s for s in prs.slides[0].shapes if s.shape_type == MSO_PICTURE]
    assert len(first_slide_pictures) >= 2

    second_slide_pictures = [s for s in prs.slides[1].shapes if s.shape_type == MSO_PICTURE]
    assert second_slide_pictures


def test_assemble_pptx_handles_template_only_slide(tmp_path: Path) -> None:
    template_path = tmp_path / "template.pptx"
    _build_template(template_path)

    pack = PackConfig(
        schema="test-pack",
        slides=[
            PackSlide(title="Static Home", id="static-home", template="single_image"),
            PackSlide(
                title="Visual Slide",
                id="visual-slide",
                template="single_image",
                visual=PackVisualRef(ref="one.yaml"),
            ),
        ],
    )

    slide_slug = slugify("visual-slide")
    png_path = tmp_path / f"{slide_slug}.png"
    _write_png(png_path)

    out = tmp_path / "deck_static.pptx"
    assemble_pack_pptx(
        pack=pack,
        results=[],
        slide_pngs={slide_slug: png_path},
        placeholder_pngs={},
        result_path=out,
        template_path=template_path,
    )

    prs = Presentation(out)
    assert len(prs.slides) == 2
    static_pictures = [s for s in prs.slides[0].shapes if s.shape_type == MSO_PICTURE]
    picture_placeholders = [
        s
        for s in prs.slides[0].shapes
        if getattr(s, "is_placeholder", False)
        and getattr(s, "placeholder_format", None)
        and s.placeholder_format.type == PP_PLACEHOLDER_ENUM.PICTURE
    ]
    assert static_pictures or picture_placeholders, "Template-only slide should retain its image without errors"


def test_assemble_pptx_adds_manual_replace_watermark(tmp_path: Path) -> None:
    template_path = tmp_path / "template.pptx"
    _build_template(template_path)

    pack = PackConfig(
        schema="test-pack",
        slides=[PackSlide(title="Manual Placeholder", id="manual-placeholder", template="single_image", manual_replace=True)],
    )

    out = tmp_path / "deck_manual_replace.pptx"
    assemble_pack_pptx(
        pack=pack,
        results=[],
        slide_pngs={},
        placeholder_pngs={},
        result_path=out,
        template_path=template_path,
    )

    prs = Presentation(out)
    assert len(prs.slides) == 1
    slide = prs.slides[0]

    replace_shapes = [
        shape
        for shape in slide.shapes
        if getattr(shape, "has_text_frame", False) and "REPLACE" in getattr(shape, "text", "")
    ]
    assert replace_shapes

    watermark = replace_shapes[0]
    assert abs(float(watermark.rotation) - 45.0) < 0.1

    runs = [run for paragraph in watermark.text_frame.paragraphs for run in paragraph.runs if run.text == "REPLACE"]
    assert runs
    assert runs[0].font.color.rgb == RGBColor(255, 0, 0)


def test_background_picture_ignored_for_single_visual_shorthand(tmp_path: Path) -> None:
    template_path = tmp_path / "template_bg.pptx"
    _build_template_with_background(template_path)

    pack = PackConfig(
        schema="test-pack",
        slides=[
            PackSlide(
                title="Home",
                id="home",
                template="home",
                visual=PackVisualRef(ref="home.yaml"),
            )
        ],
    )

    slide_slug = slugify("home")
    png_path = tmp_path / f"{slide_slug}.png"
    _write_png(png_path, width=120, height=80)

    out = tmp_path / "deck_bg.pptx"
    assemble_pack_pptx(
        pack=pack,
        results=[],
        slide_pngs={slide_slug: png_path},
        placeholder_pngs={},
        result_path=out,
        template_path=template_path,
    )

    prs = Presentation(out)
    slide = prs.slides[0]
    pictures = [shape for shape in slide.shapes if shape.shape_type == MSO_PICTURE]

    assert len(pictures) >= 2, "Background image should remain alongside replaced placeholder"
    assert "logo" in [getattr(shape, "name", None) for shape in pictures]


def test_single_visual_shorthand_preserves_slide_background_fill(tmp_path: Path) -> None:
    template_path = tmp_path / "template_bg_fill.pptx"
    _build_template_with_slide_background_fill(template_path)

    template_bg_type = Presentation(template_path).slides[0].background.fill.type
    assert template_bg_type is not None

    pack = PackConfig(
        schema="test-pack",
        slides=[
            PackSlide(
                title="Home",
                id="home",
                template="home",
                visual=PackVisualRef(ref="home.yaml"),
            )
        ],
    )

    slide_slug = slugify("home")
    png_path = tmp_path / f"{slide_slug}.png"
    _write_png(png_path, width=120, height=80)

    out = tmp_path / "deck_bg_fill.pptx"
    assemble_pack_pptx(
        pack=pack,
        results=[],
        slide_pngs={slide_slug: png_path},
        placeholder_pngs={},
        result_path=out,
        template_path=template_path,
    )

    prs = Presentation(out)
    assert len(prs.slides) == 1
    assert prs.slides[0].background.fill.type == template_bg_type
