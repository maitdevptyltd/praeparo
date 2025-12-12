from __future__ import annotations

from pptx import Presentation
from pptx.util import Inches

from praeparo.pack.formatted_values import FormattedMetricValue
from praeparo.pack.pptx import _render_jinja_in_slide
from praeparo.pack.templating import create_pack_jinja_env


def test_pptx_text_runs_render_formatted_metric_bindings() -> None:
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])

    shape = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(1))
    paragraph = shape.text_frame.paragraphs[0]
    run = paragraph.add_run()
    run.text = "We received {{ count_instructions }} instructions (raw {{ count_instructions.value }})."

    env = create_pack_jinja_env()
    context = {
        "count_instructions": FormattedMetricValue(value=54.0, format="number:0"),
    }

    _render_jinja_in_slide(slide, env, context)

    assert shape.text_frame.text == "We received 54 instructions (raw 54.0)."

