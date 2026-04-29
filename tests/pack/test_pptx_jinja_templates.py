from __future__ import annotations

from datetime import date
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches

from praeparo.models import PackConfig, PackContext, PackSlide
from praeparo.pack.pptx import (
    _iter_shapes,
    _iter_text_frames,
    _notes_template_tags,
    assemble_pack_pptx,
)


ASSET_TEMPLATE = Path(__file__).resolve().parent.parent / "assets" / "pack_template.pptx"


def _paragraph_texts(slide) -> list[str]:
    texts: list[str] = []
    for shape in _iter_shapes(slide.shapes):
        for text_frame in _iter_text_frames(shape):
            for paragraph in text_frame.paragraphs:
                texts.append(paragraph.text)
    return texts


def _build_table_template(path: Path) -> None:
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.notes_slide.notes_text_frame.text = "TEMPLATE_TAG=table_body"

    table_shape = slide.shapes.add_table(
        2,
        2,
        Inches(1),
        Inches(1),
        Inches(6),
        Inches(1.2),
    )
    table = table_shape.table
    table.cell(0, 0).text = "Month"
    table.cell(0, 1).text = "{{ display_date }}"
    table.cell(1, 0).text = "Summary"
    table.cell(1, 1).text = "Delivered for {{ team_name }}"

    prs.save(path)


def test_pptx_template_renders_display_date(tmp_path: Path) -> None:
    pack = PackConfig(
        schema="test-pack",
        context=PackContext.model_validate({"display_date": date(2025, 1, 15)}),
        slides=[PackSlide(title="Home", id="home-slide", template="home")],
    )

    out = tmp_path / "rendered.pptx"
    assemble_pack_pptx(
        pack=pack,
        results=[],
        slide_pngs={},
        placeholder_pngs={},
        result_path=out,
        template_path=ASSET_TEMPLATE,
    )

    prs = Presentation(out)
    slide = prs.slides[0]
    texts = _paragraph_texts(slide)
    joined = " ".join(texts)

    assert "{{display_date}}" not in joined
    assert "2025-01-15" in joined


def test_pptx_template_renders_jinja_inside_table_cells(tmp_path: Path) -> None:
    template = tmp_path / "table_template.pptx"
    _build_table_template(template)

    pack = PackConfig(
        schema="test-pack",
        context=PackContext.model_validate(
            {"display_date": date(2025, 1, 15), "team_name": "Operations"}
        ),
        slides=[PackSlide(title="Table Body", id="table-slide", template="table_body")],
    )

    out = tmp_path / "table_rendered.pptx"
    assemble_pack_pptx(
        pack=pack,
        results=[],
        slide_pngs={},
        placeholder_pngs={},
        result_path=out,
        template_path=template,
    )

    slide = Presentation(out).slides[0]
    table_shape = next(
        shape for shape in slide.shapes if getattr(shape, "has_table", False)
    )
    table = table_shape.table
    cell_texts = [table.cell(row, col).text for row in range(2) for col in range(2)]

    assert "{{" not in " ".join(cell_texts)
    assert "2025-01-15" in cell_texts
    assert "Delivered for Operations" in cell_texts


def test_pptx_jinja_preserves_static_text(tmp_path: Path) -> None:
    template_prs = Presentation(ASSET_TEMPLATE)
    template_lookup = _notes_template_tags(template_prs)
    expected_static = {
        "Instructions Received",
        "Documents Sent",
        "Matters Settled",
    }
    template_name = next(
        name
        for name, slide in template_lookup.items()
        if expected_static.issubset(set(_paragraph_texts(slide)))
    )

    pack = PackConfig(
        schema="test-pack",
        context=PackContext.model_validate({"display_date": date(2025, 1, 15)}),
        slides=[
            PackSlide(
                title="Highlights",
                id="highlights-slide",
                template=template_name,
            )
        ],
    )

    out = tmp_path / "static.pptx"
    assemble_pack_pptx(
        pack=pack,
        results=[],
        slide_pngs={},
        placeholder_pngs={},
        result_path=out,
        template_path=ASSET_TEMPLATE,
    )

    rendered_slide = Presentation(out).slides[0]
    rendered_static = {text for text in _paragraph_texts(rendered_slide) if text and "{{" not in text}

    assert expected_static.issubset(rendered_static)
